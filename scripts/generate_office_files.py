#!/usr/bin/env python3
"""Generate Word and PowerPoint deliverables from Finance BRD training package."""

import math
import re
import sys
from pathlib import Path

from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.shared import Inches, Pt
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches as PInches, Pt as PPt

ROOT = Path(__file__).resolve().parent.parent
OUTPUT = ROOT / "exports"
OUTPUT.mkdir(exist_ok=True)

sys.path.insert(0, str(ROOT / "ai-factory"))
from anthropic_theme import (  # noqa: E402
    SW, SH, ACCENTS, SLATE, CREAM, CLAY, SKY, OLIVE, FIG, CACTUS,
    NAVY, BLUE, GREEN, ORANGE, PURPLE, TEAL, RED, WHITE, INK, GREY, LIGHT,
    BORDER, LINK, ON_DARK_MUTED, SLATE_SOFT, CARD_FILL, MID_GRAY,
    blank as _blank, rect as _rect, text as _text, set_shape_text as _set_shape_text,
    header as _header_base, footer as _footer_base, card as _card,
    title_slide, slide_bg,
)
from anthropic_docx import (  # noqa: E402
    new_branded_document,
    add_cover,
    add_section_heading,
    add_field,
    add_body,
    add_numbered_items,
    add_callout_box,
    style_table,
)

# ======================================================================
# Shared visual-first slide library — Anthropic brand theme
# Colors: slate #141413 · cream #faf9f5 · accents clay/sky/olive
# Fonts: Poppins headings · Lora body (Arial/Georgia fallback)
# ======================================================================

FOOTER_TEXT = "Finance  ·  Internal use only"
DECK_KICKER = "Finance"


def _idx(prs):
    return len(prs.slides._sldIdLst)


def _header(slide, title, kicker=None):
    return _header_base(slide, title, kicker=kicker)


def _footer(slide, idx):
    _footer_base(slide, idx, FOOTER_TEXT)


def _frame(prs, title, kicker=None):
    """Standard content slide frame: blank + header + footer."""
    s = _blank(prs)
    y = _header(s, title, kicker=kicker)
    _footer(s, _idx(prs))
    return s, y


# ---------- text utilities ----------

_MARK = re.compile(r"^\s*([•◦\-▪☐✓✗]|\d+[.)])\s+")


def _strip_marker(line):
    return _MARK.sub("", line).strip()


def _split_kv(text):
    """Split 'Head — detail' / 'Head: detail' into (head, sub)."""
    t = text.strip()
    for sep in (" — ", " – ", " - ", ": "):
        if sep in t:
            a, b = t.split(sep, 1)
            a, b = a.strip(), b.strip()
            if a and b and len(a) <= 44:
                return a, b
    if len(t) > 60 and ", " in t:
        a, b = t.split(", ", 1)
        return a.strip(), b.strip()
    return t, None


def _ne(body):
    return [ln.rstrip() for ln in body.split("\n") if ln.strip()]


# ---------- hero / divider ----------

def _hero(prs, title, body):
    s = _blank(prs)
    title_slide(
        s,
        DECK_KICKER,
        title,
        _ne(body)[:4],
    )
    return s


# ---------- auto content renderers ----------

def _cards(s, y, items, cols, accents=None):
    """items: list of (head, sub). Renders a responsive numbered card grid."""
    n = len(items)
    rows = math.ceil(n / cols)
    gx, gy = 0.3, 0.28
    cw = (SW - 1.1 - gx * (cols - 1)) / cols
    region_h = 6.85 - y
    ch = (region_h - gy * (rows - 1)) / rows
    ch = min(ch, 2.4)
    x0 = (SW - (cw * cols + gx * (cols - 1))) / 2
    compact = ch < 1.15
    for i, (head, sub) in enumerate(items):
        r, c = divmod(i, cols)
        x = x0 + c * (cw + gx)
        yy = y + r * (ch + gy)
        acc = (accents or ACCENTS)[i % len(accents or ACCENTS)]
        _rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, yy, cw, ch, fill=CARD_FILL, line=BORDER, line_w=1)
        _rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, yy, 0.1, ch, fill=acc)
        badge = _rect(s, MSO_SHAPE.OVAL, x + 0.22, yy + 0.2, 0.42, 0.42, fill=acc)
        _set_shape_text(badge, [(str(i + 1), 13, True, CREAM)])
        if compact or not sub:
            _text(s, x + 0.78, yy, cw - 0.95, ch, [(head, 12, True, INK)],
                  anchor=MSO_ANCHOR.MIDDLE)
        else:
            _text(s, x + 0.78, yy + 0.18, cw - 0.95, 0.7, [(head, 13, True, INK)])
            _text(s, x + 0.28, yy + 0.78, cw - 0.5, ch - 0.9, [(sub, 10.5, False, GREY)])


def _checks(s, y, items, cols=2):
    n = len(items)
    rows = math.ceil(n / cols)
    gx, gy = 0.4, 0.22
    cw = (SW - 1.1 - gx * (cols - 1)) / cols
    region_h = 6.85 - y
    ch = (region_h - gy * (rows - 1)) / rows
    ch = min(ch, 1.1)
    x0 = (SW - (cw * cols + gx * (cols - 1))) / 2
    for i, text in enumerate(items):
        r, c = divmod(i, cols)
        x = x0 + c * (cw + gx)
        yy = y + r * (ch + gy)
        acc = ACCENTS[i % len(ACCENTS)]
        _rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, yy, cw, ch, fill=CARD_FILL, line=BORDER, line_w=1)
        chk = _rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, x + 0.18, yy + (ch - 0.4) / 2, 0.4, 0.4, fill=acc)
        _set_shape_text(chk, [("✓", 14, True, CREAM)])
        _text(s, x + 0.74, yy, cw - 0.9, ch, [(text, 11.5, False, INK)], anchor=MSO_ANCHOR.MIDDLE)


def _rows(s, y, pairs):
    """pairs: list of (left, right). Left navy pill -> arrow -> light box."""
    n = len(pairs)
    region_h = 6.85 - y
    gy = 0.14
    rh = (region_h - gy * (n - 1)) / n
    rh = min(rh, 0.82)
    x = 0.7
    total_w = SW - 1.4
    lw = total_w * 0.42
    rw = total_w - lw - 0.55
    for i, (left, right) in enumerate(pairs):
        yy = y + i * (rh + gy)
        acc = ACCENTS[i % len(ACCENTS)]
        lb = _rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, yy, lw, rh, fill=acc)
        _set_shape_text(lb, [(left, 11.5, True, CREAM)], align=PP_ALIGN.LEFT)
        lb.text_frame.margin_left = PInches(0.18)
        _rect(s, MSO_SHAPE.RIGHT_ARROW, x + lw + 0.06, yy + rh / 2 - 0.13, 0.4, 0.26, fill=LINK)
        rb = _rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, x + lw + 0.55, yy, rw, rh, fill=LIGHT, line=BORDER, line_w=1)
        _set_shape_text(rb, [(right, 11, False, INK)], align=PP_ALIGN.LEFT)
        rb.text_frame.margin_left = PInches(0.18)


def _two_col(s, y, left_title, left_items, left_acc, right_title, right_items, right_acc):
    cw = (SW - 1.4 - 0.4) / 2
    x0 = 0.7
    ch = 6.7 - y
    for col, (htitle, items, acc, x) in enumerate([
        (left_title, left_items, left_acc, x0),
        (right_title, right_items, right_acc, x0 + cw + 0.4),
    ]):
        _rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, cw, ch, fill=CARD_FILL, line=BORDER, line_w=1)
        _rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, cw, 0.66, fill=acc)
        _set_shape_text(_rect(s, MSO_SHAPE.RECTANGLE, x, y, cw, 0.66, fill=acc),
                        [(htitle, 15, True, CREAM)])
        iy = y + 0.95
        for it in items:
            dot = _rect(s, MSO_SHAPE.OVAL, x + 0.3, iy + 0.06, 0.16, 0.16, fill=acc)
            _text(s, x + 0.62, iy, cw - 0.85, 0.6, [(it, 11.5, False, INK)])
            iy += max(0.5, 0.42 + 0.18 * (len(it) // 46))


def _callout(prs, s, y, title, body):
    lines = _ne(body)
    box_y = y + 0.4
    box_h = 6.6 - box_y
    _rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, 1.4, box_y, SW - 2.8, box_h, fill=LIGHT, line=SLATE, line_w=1.5)
    _rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, 1.4, box_y, 0.14, box_h, fill=CLAY)
    yy = box_y + (box_h - len(lines) * 0.62) / 2
    for ln in lines:
        emph = any(tok in ln for tok in ("[", "]")) or ln.isupper()
        size = 22 if (len(ln) < 60 and emph) else 17
        _text(s, 1.9, yy, SW - 3.8, 0.7, [(ln, size, emph, SLATE if emph else INK)],
              align=PP_ALIGN.CENTER)
        yy += 0.62


def _chevrons(s, y, steps, note=None):
    n = len(steps)
    gap = 0.12
    cw = (SW - 1.0 - gap * (n - 1)) / n
    ch = 1.7
    x0 = 0.5
    yy = y + 0.7
    for i, label in enumerate(steps):
        acc = ACCENTS[i % len(ACCENTS)]
        sp = _rect(s, MSO_SHAPE.CHEVRON, x0 + i * (cw + gap), yy, cw, ch, fill=acc, line=CREAM, line_w=1.5)
        _set_shape_text(sp, [(str(i + 1), 16, True, CREAM), (label, 11.5, True, CREAM)])
    if note:
        _text(s, 0.55, yy + ch + 0.5, 12.2, 0.8, [(note, 13, False, GREY)], align=PP_ALIGN.CENTER)


def _line_tiles(s, y, lines):
    items = [(_split_kv(ln)) for ln in lines]
    cols = 1 if len(items) <= 3 else 2
    _cards(s, y, items, cols)


# ---------- dispatcher ----------

def _looks_mapping(lines):
    hits = 0
    for ln in lines:
        if "→" in ln and ln.count("→") == 1:
            hits += 1
        elif re.search(r"\S\s{2,}\S", ln) and "•" not in ln:
            hits += 1
    return hits >= max(2, int(len(lines) * 0.6))


def _to_pair(ln):
    if "→" in ln:
        a, b = ln.split("→", 1)
    else:
        a, b = re.split(r"\s{2,}", ln.strip(), maxsplit=1)
    return a.strip(), b.strip()


def _render_visual(prs, title, body):
    """Convert a (title, body) content tuple into a visual slide."""
    if _idx(prs) == 0:
        return _hero(prs, title, body)

    lines = _ne(body)
    raw = [ln for ln in lines]
    flat = "\n".join(lines)

    # 1) Multi-arrow single flow -> chevrons
    arrow_lines = [ln for ln in lines if "→" in ln]
    total_arrows = flat.count("→")
    if total_arrows >= 2 and len(arrow_lines) <= 2 and len(lines) <= 3:
        seq = re.split(r"→|\n", flat)
        steps = [s.strip() for s in seq if s.strip()]
        if 3 <= len(steps) <= 8:
            s, y = _frame(prs, title)
            _chevrons(s, y, steps)
            return s

    # 2) Comparison (good vs bad / weak vs strong / do vs don't)
    up = flat.upper()
    cmp_pairs = [("BAD", "GOOD"), ("WEAK", "STRONG"), ("DON'T", "DO")]
    for neg, pos in cmp_pairs:
        if neg in up and pos in up:
            s, y = _frame(prs, title)
            neg_items, pos_items, bucket = [], [], None
            for ln in lines:
                u = ln.upper().strip(' :"')
                if u.startswith(neg):
                    bucket = "neg"
                    rest = ln.split(":", 1)[1].strip() if ":" in ln else ""
                    if rest:
                        neg_items.append(rest.strip('"'))
                    continue
                if u.startswith(pos):
                    bucket = "pos"
                    rest = ln.split(":", 1)[1].strip() if ":" in ln else ""
                    if rest:
                        pos_items.append(rest.strip('"'))
                    continue
                t = _strip_marker(ln).strip('"')
                if bucket == "neg":
                    neg_items.append(t)
                elif bucket == "pos":
                    pos_items.append(t)
            if neg_items and pos_items:
                _two_col(s, y, f"✗  {neg.title()}", neg_items, RED,
                         f"✓  {pos.title()}", pos_items, GREEN)
                return s

    # ✓ / ✗ comparison
    do_items = [_strip_marker(ln) for ln in lines if ln.lstrip().startswith("✓")]
    dont_items = [_strip_marker(ln) for ln in lines if ln.lstrip().startswith("✗")]
    if do_items and dont_items:
        s, y = _frame(prs, title)
        _two_col(s, y, "✓  Do include", do_items, GREEN, "✗  Do NOT include", dont_items, RED)
        return s

    # 3) Checklist
    check_lines = [ln for ln in lines if ln.lstrip().startswith("☐")]
    if len(check_lines) >= 3 and len(check_lines) >= len(lines) - 2:
        s, y = _frame(prs, title)
        caption = next((ln for ln in lines if not ln.lstrip().startswith("☐")), None)
        if caption and len(caption) <= 90:
            _text(s, 0.7, y, 12.0, 0.4, [(caption, 13, False, GREY)])
            y += 0.55
        items = [_strip_marker(ln) for ln in check_lines]
        _checks(s, y, items, cols=2 if len(items) <= 10 else 2)
        return s

    # 4) Numbered list
    num_lines = [ln for ln in lines if re.match(r"^\s*\d+[.)]\s", ln)]
    if len(num_lines) >= 3 and len(num_lines) >= len(lines) - 2:
        s, y = _frame(prs, title)
        items = [_split_kv(_strip_marker(ln)) for ln in num_lines]
        n = len(items)
        cols = 2 if n <= 6 else (3 if n <= 9 else 2)
        _cards(s, y, items, cols)
        return s

    # 5) Mapping rows
    map_candidates = [ln for ln in lines if "•" not in ln and not ln.lstrip().startswith("☐")]
    if len(map_candidates) >= 3 and _looks_mapping(map_candidates):
        s, y = _frame(prs, title)
        pairs = [_to_pair(ln) for ln in map_candidates if "→" in ln or re.search(r"\S\s{2,}\S", ln)]
        if len(pairs) >= 3:
            _rows(s, y, pairs[:8])
            return s

    # 6) Bullet cards
    bullets = [ln for ln in lines if ln.lstrip().startswith(("•", "◦", "▪", "-"))]
    if len(bullets) >= 2:
        s, y = _frame(prs, title)
        lead = lines[0] if not lines[0].lstrip().startswith(("•", "◦", "▪", "-")) else None
        if lead and len(lead) <= 90:
            _text(s, 0.7, y, 12.0, 0.4, [(lead, 13, False, GREY)])
            y += 0.5
        items = [_split_kv(_strip_marker(ln)) for ln in bullets]
        n = len(items)
        cols = 1 if n <= 2 else (2 if n <= 6 else (3 if n <= 9 else 2))
        _cards(s, y, items, cols)
        return s

    # 7) Short callout (formula / golden rule / definition)
    if len(lines) <= 6:
        s, y = _frame(prs, title)
        _callout(prs, s, y, title, body)
        return s

    # 8) Fallback: line tiles
    s, y = _frame(prs, title)
    _line_tiles(s, y, [_strip_marker(ln) for ln in lines[:8]])
    return s


def add_heading(doc, text, level=1):
    """Legacy wrapper — use add_section_heading for new code."""
    add_section_heading(doc, text, accent_index=max(0, level - 1))


def generate_brd_template_docx():
    doc = new_branded_document()
    add_cover(
        doc,
        "Finance — Business Requirements Document (BRD)",
        meta="Template v1.0 · Internal use only",
    )
    add_callout_box(
        doc,
        "Instructions: Complete all mandatory sections. Write what the business needs and why — "
        "not how IT should build it. Obtain Business Sponsor sign-off before IT/BA intake.",
    )

    sections = [
        ("A. REQUEST INFORMATION", [
            "BRD Title:", "Business Unit:", "Requester Name:", "Requester Email:",
            "Business Sponsor (Director+):", "Date Submitted:", "Target Go-Live Date:",
            "Priority Category:", "Regulatory Deadline (if applicable):",
        ]),
        ("B. EXECUTIVE SUMMARY", [
            "Business problem:", "Who is impacted:", "Business impact if not addressed:",
            "Proposed business outcome (one sentence):",
        ]),
        ("C. BUSINESS OBJECTIVES & SUCCESS METRICS", [
            "Objective 1 — Baseline / Target / Measurement:",
            "Objective 2 — Baseline / Target / Measurement:",
            "Objective 3 — Baseline / Target / Measurement:",
        ]),
        ("D. BACKGROUND & CURRENT STATE", [
            "Current process (step-by-step):", "Current systems involved:",
            "Volume and frequency:",
        ]),
        ("E. PROPOSED TO-BE STATE", [
            "Desired future process (business view only):", "Expected changes:",
        ]),
        ("F. SCOPE", ["In scope:", "Out of scope:", "Assumptions:", "Dependencies:"]),
        ("G. STAKEHOLDERS & USERS", ["Role / Department / # users / Location / Access:"]),
        ("H. BUSINESS RULES", [
            "Eligibility rules:", "Approval authority:", "Pricing/fees (business level):",
            "Exception handling:", "Cut-off times:", "Customer communication rules:",
        ]),
        ("I. DATA & SECURITY REQUIREMENTS", [
            "Data types:", "Classification:", "Data sharing:", "Retention:",
            "Audit requirements:", "Customer consent:", "Remote access implications:",
        ]),
        ("J. REPORTING & CONTROLS", [
            "Reports required:", "Control points:", "Reconciliation:", "Fraud/risk controls:",
        ]),
        ("K. IMPLEMENTATION & ROLLOUT", [
            "Rollout approach:", "Training needs:", "Language:", "Business continuity fallback:",
        ]),
        ("L. RISKS", ["Operational / Customer / Financial / Compliance risks:"]),
        ("M. ACCEPTANCE CRITERIA (min 5)", [
            "AC-01: Given / When / Then",
            "AC-02: Given / When / Then",
            "AC-03: Given / When / Then",
            "AC-04: Given / When / Then",
            "AC-05: Given / When / Then",
        ]),
        ("N. COMPLIANCE SCREENING", [
            "Changes origination/disbursement? Y/N",
            "Changes interest/fees/contract? Y/N",
            "Sends customer notifications? Y/N",
            "Third-party data exchange? Y/N",
            "Affects collections/legal? Y/N",
            "Audit trail required? Y/N",
            "Exposes data on POS/field devices? Y/N",
        ]),
        ("O. APPROVALS", [
            "Business Requester:", "Business Sponsor:", "Business Unit Head:",
            "Risk/Compliance (if required):", "Data Owner (if required):",
        ]),
    ]

    for i, (title, fields) in enumerate(sections):
        add_section_heading(doc, title, accent_index=i)
        for field in fields:
            label = field if field.endswith(":") else f"{field}:"
            add_field(doc, label)

    out = OUTPUT / "Finance-BRD-Template.docx"
    doc.save(out)
    print(f"Created {out}")
    return out


def generate_cheat_sheet_docx():
    doc = new_branded_document()
    add_cover(
        doc,
        "Finance BRD Cheat Sheet",
        subtitle="Quick reference for business authors · POS, HQ, Digital, Collections",
        meta="Internal use only",
    )
    add_callout_box(
        doc,
        "Golden rule: write the problem and business rules — not the system design. "
        "Aim for ≥ 80% on the BRD quality gate before IT intake.",
    )

    for i, (lang_title, items) in enumerate([
        ("Tiếng Việt (POS / hiện trường)", [
            "Bắt đầu bằng vấn đề — không bắt đầu bằng hệ thống",
            "Có số liệu: %, giờ, số KH",
            "Ghi rõ ai dùng: vai trò, số người, địa điểm",
            "Quy tắc nghiệp vụ: ai duyệt, giờ cắt, ngoại lệ",
            "Trong phạm vi + ngoài phạm vi",
            "Dữ liệu & tuân thủ: PII, CIC, SMS",
            "Tiêu chí nghiệm thu: Cho trước / Khi / Thì (≥5)",
            "Xin Sponsor (Giám đốc+) ký trước khi gửi IT",
        ]),
        ("English (HQ / Digital / Collections)", [
            "Start with the problem, not the system",
            "Quantify impact: %, hours, customers, VND",
            "Define who uses it: role, count, location",
            "Document business rules: approval, cut-off, exceptions",
            "Write in-scope AND out-of-scope",
            "Complete data & compliance section honestly",
            "Add acceptance criteria: Given/When/Then (min 5)",
            "Get Sponsor sign-off before IT intake",
        ]),
    ]):
        add_section_heading(doc, lang_title, accent_index=i)
        add_numbered_items(doc, items)

    add_section_heading(doc, "Compliance triggers → Route to", accent_index=2)
    table = doc.add_table(rows=1, cols=2)
    hdr = table.rows[0].cells
    hdr[0].text = "If BRD includes..."
    hdr[1].text = "Route to"
    style_table(table, ["If BRD includes...", "Route to"])
    rows = [
        ("Loan approval / disbursement change", "Risk + Credit Policy"),
        ("Interest / fees / contract change", "Legal + Finance"),
        ("Customer SMS / email", "Legal (template)"),
        ("Third-party data sharing", "Vendor Risk + Legal"),
        ("Collections / legal process", "Compliance"),
        ("Data on POS / home device", "Security"),
    ]
    for a, b in rows:
        row = table.add_row().cells
        row[0].text = a
        row[1].text = b

    out = OUTPUT / "Finance-BRD-Cheat-Sheet.docx"
    doc.save(out)
    print(f"Created {out}")
    return out


SLIDES = [
    ("Finance BRD Training", "Write requirements IT can deliver — without rework\nSession 1 of 4"),
    ("Learning Objectives", "• Write a complete BRD using Finance template\n• Separate business need from technical solution\n• Define testable acceptance criteria\n• Pass BRD quality gate (≥ 80%)"),
    ("Why We Are Here", "• High volume of IT requests from business\n• Rework costs: delay, UAT defects, compliance risk\n• Goal: faster triage, predictable delivery"),
    ("What Is a BRD?", "Document describing:\n• Business problem\n• Outcomes and KPIs\n• Business rules and scope\n\nNOT: technical design, vendor selection, project plan"),
    ("BRD vs Other Documents", "BRD → Business owns\nFRD → BA + IT\nTechnical Design → IT\nUAT → Business + QA"),
    ("Finance Context", "• FE ONLINE 2.0 (customer mobile app)\n• Finacle LMS, CIF, Assure on AWS\n• POS/LOS, FirstVision cards, Collections\n• ~13k POS, SBV compliance expectations"),
    ("Request Journey", "Business BRD → BA review → IT triage → Risk (if needed)\n→ FRD → Build → UAT → Go-live"),
    ("Good vs Bad Request", "BAD: 'Build Kafka API to Finacle'\nGOOD: 'POS needs same-day approval visibility during customer visit'"),
    ("The Golden Rule", "Write the PROBLEM and RULES,\nnot the SYSTEM DESIGN."),
    ("Problem Statement Formula", "[Role] cannot [task] because [constraint],\ncausing [impact].\n\nExample: POS staff cannot confirm approved installment plan during customer visit, causing 18% drop-off."),
    ("Objectives Must Be Measurable", "WEAK: Improve customer experience\nSTRONG: Reduce disbursement TAT from 24h to 4h for POS motorbike loans"),
    ("Scope: In AND Out", "In scope: what this project delivers\nOut of scope: prevents scope creep\n\nAlways write both."),
    ("Business Rules — IT Must Not Guess", "• Eligibility\n• Approval authority\n• Cut-off times\n• Exceptions\n• Customer communication"),
    ("Acceptance Criteria", "Given [context]\nWhen [action]\nThen [expected result]\n\nMinimum 5 per BRD → becomes UAT tests"),
    ("Data & Compliance", "• Classify data: Public / Internal / Confidential / Restricted\n• Complete Section N screening\n• Route to Risk if customer data on POS or notifications"),
    ("Top 10 BRD Mistakes", "1. Solution not need\n2. No KPI baseline\n3. Missing business rules\n4. No out-of-scope\n5. Multiple projects in one BRD\n6. No acceptance criteria\n7. No sponsor\n8. Ignoring POS training\n9. No data classification\n10. Vendor brochure copy-paste"),
    ("BRD Quality Gate", "• Score ≥ 80% on quality rubric\n• Sponsor sign-off\n• All mandatory sections complete\n• Then → IT triage"),
    ("Certification: BRD Ready", "• Attend all 4 sessions\n• Submit 1 BRD scoring ≥ 80%\n• Sponsor sign-off\n• BA coach confirms readiness"),
    ("Next Steps", "• Use template on Confluence\n• Office hours: weekly BA support\n• Submit via ServiceNow/Jira BRD form\n• Questions: BA team"),
]


def _add_content_slide(prs, title, body):
    return _render_visual(prs, title, body)


def _add_pipeline_slide(prs, title, stages, subtitle=None):
    """Horizontal pipeline diagram with rounded boxes (restyled)."""
    s, y = _frame(prs, title)
    if subtitle:
        _text(s, 0.7, y, 12.0, 0.4, [(subtitle, 13, False, GREY)])

    n = len(stages)
    total_w = SW - 1.0
    box_w = min(1.7, (total_w - 0.3 * (n - 1)) / n)
    gap = 0.3
    start_x = (SW - (box_w * n + gap * (n - 1))) / 2
    yb = 3.4
    h = 1.05
    for i, label in enumerate(stages):
        x = start_x + i * (box_w + gap)
        acc = ACCENTS[i % len(ACCENTS)]
        box = _rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, yb, box_w, h, fill=acc, line=CREAM, line_w=1.5)
        _set_shape_text(box, [(label, 11, True, CREAM)])
        if i < n - 1:
            _rect(s, MSO_SHAPE.RIGHT_ARROW, x + box_w + 0.02, yb + h / 2 - 0.16, gap - 0.02, 0.32, fill=LINK)
    return s


def _add_dual_track_slide(prs):
    """Governance track (top) + Agile delivery track (bottom), restyled."""
    s, y = _frame(prs, "Dual-Track: Governance + Agile Delivery")

    tracks = [
        ("Governance & assurance — runs in parallel", [
            "BRD gate", "GRC / Legal", "IT-Security", "ARB", "CAB", "Audit evidence",
        ], 2.1, FIG),
        ("Agile delivery — Scrum (2-week sprints)", [
            "FSD", "Sprint plan", "Build", "SIT", "UAT", "Ship", "Hypercare",
        ], 4.0, OLIVE),
    ]
    n_boxes = 7
    box_w = 1.55
    gap = 0.22
    start_x = (SW - (box_w * n_boxes + gap * (n_boxes - 1))) / 2

    for track_name, stages, y_in, accent in tracks:
        _text(s, 0.5, y_in - 0.42, 9.0, 0.35, [(track_name, 13, True, accent)])
        for i, label in enumerate(stages):
            x = start_x + i * (box_w + gap)
            box = _rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, y_in, box_w, 0.78, fill=accent, line=CREAM, line_w=1.5)
            _set_shape_text(box, [(label, 10, True, CREAM)])

    _text(s, 0.7, 5.35, 12.0, 0.5,
          [("Governance runs in parallel — it does not wait for sprint end.", 12.5, False, INK)])
    _text(s, 0.7, 5.8, 12.0, 0.5,
          [("Ship only when BOTH tracks are green.", 12.5, True, NAVY)])
    return s


def _add_scrum_cycle_slide(prs):
    """Scrum ceremonies mapped to delivery phases — three styled phase columns."""
    s, y = _frame(prs, "Scrum Ceremonies Inside Delivery")
    phases = [
        (BLUE, "Sprint N", "≈ 2 weeks", [
            ("Mon", "Sprint Planning", "Pull stories from FSD"),
            ("Gate", "Entry check", "BRD + FSD linked, flags clear"),
            ("Daily", "Stand-up 15 min", "Dev + SM + Ops"),
            ("Wed", "Mid-sprint", "BA clarification, no scope creep"),
            ("Fri", "Code complete", "Unit tests + SIT entry"),
        ]),
        (ORANGE, "Sprint N+1", "Release sprint", [
            ("Mon", "SIT / QA", "Regression, exit criteria"),
            ("Wed", "UAT window", "Business tests BRD AC"),
            ("Thu", "Sprint Review", "Demo to Sponsor delegate"),
            ("Fri", "Retrospective", "Process improvement"),
        ]),
        (GREEN, "Release week", "Sprint boundary", [
            ("Prep", "Release readiness", "Ops checklist + CAB"),
            ("Go", "Ship to prod", "Deploy + verify"),
            ("T+1", "Hypercare", "Handover to support"),
        ]),
    ]
    cw, gap = 4.0, 0.3
    x0 = (SW - (cw * 3 + gap * 2)) / 2
    ch = 4.9
    for i, (acc, name, when, rows) in enumerate(phases):
        x = x0 + i * (cw + gap)
        _rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, y + 0.2, cw, ch, fill=CARD_FILL, line=BORDER, line_w=1)
        _set_shape_text(_rect(s, MSO_SHAPE.RECTANGLE, x, y + 0.2, cw, 0.75, fill=acc),
                        [(name, 15, True, CREAM), (when, 10, False, ON_DARK_MUTED)])
        ry = y + 1.15
        for tag, head, detail in rows:
            pill = _rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, x + 0.2, ry, 0.85, 0.5, fill=LIGHT, line=acc, line_w=1)
            _set_shape_text(pill, [(tag, 9.5, True, acc)])
            _text(s, x + 1.15, ry - 0.04, cw - 1.3, 0.3, [(head, 11, True, INK)])
            _text(s, x + 1.15, ry + 0.26, cw - 1.3, 0.3, [(detail, 9, False, GREY)])
            ry += 0.72
    return s


def _add_swimlane_slide(prs, title, columns, lanes, subtitle=None):
    """Cross-functional swimlane: role lanes (rows) x phase columns.

    lanes: list of (lane_name, [cell_text_or_None per column]).
    """
    s, top = _frame(prs, title)
    if subtitle:
        _text(s, 0.7, top, 12.0, 0.35, [(subtitle, 13, False, GREY)])
        top += 0.5

    label_w = 1.7
    grid_x = 0.4 + label_w
    grid_w = SW - grid_x - 0.25
    n_cols = len(columns)
    col_w = grid_w / n_cols

    for j, col in enumerate(columns):
        _text(s, grid_x + j * col_w, top, col_w, 0.4, [(col, 11, True, NAVY)], align=PP_ALIGN.CENTER)

    lane_top = top + 0.45
    lane_h = (6.9 - lane_top) / len(lanes)
    pad = 0.06

    for i, (lane_name, cells) in enumerate(lanes):
        accent = ACCENTS[i % len(ACCENTS)]
        y = lane_top + i * lane_h
        lbl = _rect(s, MSO_SHAPE.RECTANGLE, 0.4, y, label_w, lane_h - pad, fill=accent, line=CREAM, line_w=1)
        _set_shape_text(lbl, [(lane_name, 10.5, True, CREAM)])
        _rect(s, MSO_SHAPE.RECTANGLE, grid_x, y, grid_w, lane_h - pad, fill=LIGHT, line=BORDER, line_w=1)
        for j, cell in enumerate(cells):
            if not cell:
                continue
            box = _rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, grid_x + j * col_w + pad, y + pad,
                        col_w - 2 * pad, lane_h - pad - 2 * pad, fill=accent, line=CREAM, line_w=1)
            _set_shape_text(box, [(cell, 9, True, CREAM)])

    return s


def _add_flowchart_slide(prs, title, steps, subtitle=None):
    """Vertical decision flowchart.

    steps: list of dicts with keys:
      kind: "start" | "process" | "decision" | "end"
      text: label
      branch: optional (label, outcome_text) drawn to the right (decision only)
    """
    s, top = _frame(prs, title)
    if subtitle:
        _text(s, 0.7, top, 12.0, 0.35, [(subtitle, 13, False, GREY)])
        top += 0.45

    n = len(steps)
    avail = 6.9 - top
    row_h = avail / n
    node_h = min(0.85, row_h - 0.18)
    main_x = 1.0
    main_w = 4.6
    out_x = 6.6
    out_w = 6.0

    kind_color = {"start": OLIVE, "process": SKY, "decision": CLAY, "end": FIG}
    kind_shape = {
        "start": MSO_SHAPE.ROUNDED_RECTANGLE,
        "process": MSO_SHAPE.RECTANGLE,
        "decision": MSO_SHAPE.DIAMOND,
        "end": MSO_SHAPE.ROUNDED_RECTANGLE,
    }

    for i, step in enumerate(steps):
        y = top + i * row_h
        kind = step["kind"]
        shape = _rect(s, kind_shape[kind], main_x, y, main_w, node_h,
                      fill=kind_color[kind], line=CREAM, line_w=1.5)
        _set_shape_text(shape, [(step["text"], 10, True, CREAM)])

        if i < n - 1:
            conn = s.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT,
                PInches(main_x + main_w / 2), PInches(y + node_h),
                PInches(main_x + main_w / 2), PInches(top + (i + 1) * row_h),
            )
            conn.line.color.rgb = LINK
            conn.line.width = PPt(1.5)

        branch = step.get("branch")
        if branch:
            blabel, btext = branch
            ob = _rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, out_x, y + (node_h - 0.6) / 2, out_w, 0.6,
                       fill=SKY, line=CREAM, line_w=1.5)
            _set_shape_text(ob, [(btext, 9.5, True, CREAM)])
            conn = s.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT,
                PInches(main_x + main_w), PInches(y + node_h / 2),
                PInches(out_x), PInches(y + node_h / 2),
            )
            conn.line.color.rgb = LINK
            conn.line.width = PPt(1.5)
            _text(s, main_x + main_w + 0.05, y + node_h / 2 - 0.34,
                  out_x - main_x - main_w, 0.3, [(blabel, 9, True, GREEN)], align=PP_ALIGN.CENTER)

    return s


def _add_org_chart_slide(prs, title, root, groups, subtitle=None):
    """Org chart: one root box, then manager columns each with stacked reports.

    groups: list of (manager_label, [report_labels]).
    """
    s, top = _frame(prs, title)
    if subtitle:
        _text(s, 0.7, top, 12.0, 0.35, [(subtitle, 12.5, False, GREY)])
        top += 0.5

    root_w, root_h = 5.0, 0.72
    root_x = (SW - root_w) / 2
    rb = _rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, root_x, top, root_w, root_h, fill=SLATE, line=CREAM, line_w=1.5)
    _set_shape_text(rb, [(root, 13, True, CREAM)])

    n = len(groups)
    col_gap = 0.3
    col_w = (SW - 0.6 - col_gap * (n - 1)) / n
    start_x = 0.3
    mgr_top = top + root_h + 0.6
    mgr_h = 0.68
    report_h = 0.58
    report_gap = 0.18

    for i, (mgr, reports) in enumerate(groups):
        accent = ACCENTS[i % len(ACCENTS)]
        cx = start_x + i * (col_w + col_gap)
        mgr_cx = cx + col_w / 2
        conn = s.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            PInches(root_x + root_w / 2), PInches(top + root_h),
            PInches(mgr_cx), PInches(mgr_top),
        )
        conn.line.color.rgb = LINK
        conn.line.width = PPt(1.5)
        mb = _rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, cx, mgr_top, col_w, mgr_h, fill=accent, line=CREAM, line_w=1.5)
        _set_shape_text(mb, [(mgr, 11, True, CREAM)])
        ry = mgr_top + mgr_h + 0.35
        for report in reports:
            rbx = _rect(s, MSO_SHAPE.RECTANGLE, cx + 0.2, ry, col_w - 0.4, report_h, fill=LIGHT, line=accent, line_w=1)
            _set_shape_text(rbx, [(report, 9.5, False, INK)])
            ry += report_h + report_gap

    return s


IT_DELIVERY_SLIDES = [
    (
        "Finance IT Delivery Framework",
        "Requirement → BRD → FSD → Test → UAT → Ship → Operations Support\n"
        "Governed delivery with Agile/Scrum internal process",
    ),
    (
        "Learning Objectives",
        "• Map the end-to-end IT delivery pipeline and mandatory gates\n"
        "• Align Scrum ceremonies with BRD / FSD / release governance\n"
        "• Clarify roles: Business, Dev, IT Ops, IT-Governance, IT-Security, GRC\n"
        "• Apply evidence chain from requirement intake to operations support",
    ),
    (
        "Document Chain — Who Owns What",
        "BRD (Business Requirements)     → Business writes, Sponsor signs\n"
        "FSD / FRD (Functional Spec)     → BA + IT Product; derived from BRD\n"
        "User Stories & Tasks            → BA + Dev; traceable to FSD + BRD AC\n"
        "Test Cases (SIT)                → QA; mapped to acceptance criteria\n"
        "UAT Scripts                     → Business; from BRD Section M\n"
        "Change Request (CR)             → IT Ops + IT-Governance (CAB)\n\n"
        "Rule: each document links upward — full traceability for Audit",
    ),
    (
        "Three Lanes Operating Model",
        "LANE 1 — DEMAND (Business + BA)\n"
        "  Define need → BRD → Sponsor → BA quality ≥80% → compliance routing → triage\n\n"
        "LANE 2 — DELIVERY (IT Product + Dev + Ops)\n"
        "  FSD → Scrum sprints → SIT → UAT support → CAB → Ship → Hypercare\n\n"
        "LANE 3 — ASSURANCE (IT-Gov + IT-Sec + GRC + Audit)\n"
        "  Architecture standards, security review, regulatory control, evidence pack",
    ),
    (
        "Stakeholder Map",
        "BUSINESS          Sponsor, requesters, UAT owners\n"
        "IT DELIVERY       Ops Manager, BA, IT Product, Dev, Service Desk\n"
        "IT-GOVERNANCE     CAB, ARB, CMDB, intake policy, change control\n"
        "IT-SECURITY       Security architecture, IAM, exceptions, secure SDLC\n"
        "GRC / LEGAL       Business compliance, contracts, collections\n"
        "AUDIT / BOD       Evidence, material risk reporting",
    ),
    (
        "Routing Rule — IT-Gov vs IT-Sec vs GRC",
        "IT process & change control     → IT-Governance\n"
        "Security, access, data protection → IT-Security\n"
        "Business & regulatory controls  → Risk / GRC / Legal\n\n"
        "All three may apply to one BRD — route early, run in parallel",
    ),
    (
        "Phase 1 — Receive Requirement",
        "Channel: ServiceNow / Jira BRD catalog (not email for projects)\n\n"
        "1. Classify: standard IT ticket vs business change vs build request\n"
        "2. Wrong bucket? → redirect to BRD (business must define first)\n"
        "3. Assign: Service Desk | BA | IT Product\n\n"
        "Ops Manager: enforce one front door; log informal requests",
    ),
    (
        "Phase 2 — BRD & Quality Gate",
        "Business drafts BRD → Sponsor (Director+) signs\n"
        "BA quality review (≤2 business days)\n"
        "  • Score ≥ 80% on quality rubric\n"
        "  • Mandatory sections complete\n"
        "  • Compliance Section N → auto-route GRC / IT-Security\n"
        "Outcome: Accepted to triage | Returned | Deferred",
    ),
    (
        "Phase 3 — IT Triage & Product Backlog",
        "IT Product / Portfolio (≤5 business days)\n"
        "  • Fit with strategy and capacity\n"
        "  • Priority vs regulatory / revenue / risk\n"
        "  • Decision: Approved for FSD | Deferred | Rejected (with reason)\n\n"
        "Approved BRD → Product Backlog item (Epic) in Jira",
    ),
    (
        "Phase 4 — FSD (Functional Specification / FRD)",
        "BA + IT Product create FSD from accepted BRD\n"
        "FSD contains:\n"
        "  • Functional flows mapped to BRD to-be state\n"
        "  • Interface & system touchpoints (FE ONLINE, Finacle, LOS…)\n"
        "  • Story breakdown with traceability ID → BRD-xxx\n"
        "  • Non-functional: performance, security, audit trail\n\n"
        "Gate: no sprint planning without FSD linked to BRD",
    ),
    (
        "Agile / Scrum — Where It Fits",
        "Governance gates are OUTSIDE the sprint (BRD, FSD, CAB).\n"
        "Scrum runs INSIDE build & test phases.\n\n"
        "Epic     = Accepted BRD\n"
        "Feature  = FSD capability group\n"
        "Story    = Implementable unit with AC\n"
        "Task     = Dev work item\n\n"
        "2-week sprints typical; release may span 1–2 sprints + UAT window",
    ),
    (
        "Scrum Roles at Finance",
        "Product Owner (IT Product)    Prioritizes backlog; accepts sprint output\n"
        "Scrum Master (Dev Lead)       Facilitates ceremonies; removes blockers\n"
        "Development Team            Build, unit test, support SIT fixes\n"
        "BA                          FSD, story refinement, UAT support\n"
        "IT Ops Manager                Release discipline, CAB, hypercare\n"
        "Business Sponsor              UAT accountability; go-live sign-off",
    ),
    (
        "Backlog Hierarchy",
        "Product Backlog (prioritized by IT Product)\n"
        "  └── Epic [BRD-2026-0042] POS lending visibility\n"
        "        └── Feature [FSD-3.2] Approval status API to POS\n"
        "              └── Story [US-101] Display pre-approved amount\n"
        "              └── Story [US-102] Handle expired pre-approval\n"
        "                    └── Tasks: dev, test, docs\n\n"
        "Every story must reference BRD acceptance criteria IDs",
    ),
    (
        "Sprint Planning — Entry Checklist",
        "☐ BRD accepted (score ≥ 80%) and linked\n"
        "☐ FSD approved for this epic / feature\n"
        "☐ GRC / Legal review complete (if flagged)\n"
        "☐ IT-Security review complete (if Restricted / integration)\n"
        "☐ Stories have clear AC and testable outcomes\n"
        "☐ Team capacity confirmed; no conflicting release\n"
        "☐ Ops aware of target UAT / deploy window",
    ),
    (
        "Definition of Done — Development Team",
        "☐ Code complete + peer review\n"
        "☐ Unit tests pass; coverage per team standard\n"
        "☐ Traceability: story → BRD AC → test case\n"
        "☐ SIT-ready build deployed to test environment\n"
        "☐ Release notes draft for Ops\n"
        "☐ No open Critical/High defects for in-scope stories\n"
        "☐ Security scan clean or exception recorded (IT-Security)",
    ),
    (
        "Phase 5 — Test (SIT / QA)",
        "QA executes SIT against FSD + BRD acceptance criteria\n\n"
        "Flow: Unit tests (Dev) → SIT (QA) → Regression → Defect fix loop\n\n"
        "Exit criteria:\n"
        "  • All in-scope test cases Pass\n"
        "  • No open Sev-1 / Sev-2 defects\n"
        "  • Test summary report attached to release record\n\n"
        "Then → UAT environment ready for Business",
    ),
    (
        "Phase 6 — UAT (Business Validation)",
        "Business executes UAT scripts derived from BRD Section M\n"
        "  Given / When / Then — objective pass/fail\n\n"
        "IT provides: UAT environment, test data, defect fixes\n"
        "Business provides: testers, sign-off, edge-case validation\n\n"
        "Gate: Sponsor (or delegate) signs UAT record\n"
        "IT does NOT accept on behalf of business",
    ),
    (
        "Phase 7 — Ship (Release to Production)",
        "1. Release readiness review (Ops checklist)\n"
        "2. Change Request submitted — linked BRD, FSD, UAT sign-off\n"
        "3. CAB approval (IT-Governance) — IT-Security if flagged\n"
        "4. Deploy via approved pipeline / runbook\n"
        "5. Post-deploy verification (smoke test + monitoring)\n"
        "6. CMDB / release record updated\n\n"
        "Hard gate: no production deploy without UAT + CAB",
    ),
    (
        "Phase 8 — Operations Support",
        "HYPERCARE (T+1 to T+14)\n"
        "  • Ops war room for Critical releases\n"
        "  • Dev on standby for defects\n"
        "  • Daily check with Business Sponsor\n\n"
        "BAU HANDOVER\n"
        "  • Service Desk knowledge article\n"
        "  • Monitoring & alerting configured\n"
        "  • Incident routing defined\n"
        "  • Evidence pack closed for Audit (13 artifacts)",
    ),
    (
        "Release Readiness — Go / No-Go",
        "☐ BRD accepted and linked\n"
        "☐ FSD complete; stories closed\n"
        "☐ SIT exit criteria met\n"
        "☐ UAT signed by Business Sponsor\n"
        "☐ CAB approved (IT-Governance)\n"
        "☐ IT-Security sign-off (if required)\n"
        "☐ Rollback plan tested\n"
        "☐ Comms & training complete (BRD Section K)\n"
        "☐ Hypercare roster assigned",
    ),
    (
        "Hard Gates — Non-Negotiable",
        "1. No FSD / sprint work without accepted BRD (≥ 80%)\n"
        "2. No story without traceability to BRD acceptance criteria\n"
        "3. No UAT skip — business signs, not IT\n"
        "4. No production deploy without CAB (IT-Governance)\n"
        "5. No security bypass without IT-Security exception register\n"
        "6. Email / chat is not intake for projects",
    ),
    (
        "Audit Evidence Pack (per release)",
        "1. BRD (+ versions)  2. BA scorecard  3. Sponsor approval\n"
        "4. IT-Security sign-off (if flagged)  5. GRC/Legal (if flagged)\n"
        "6. FSD linked to BRD  7. SIT results  8. UAT sign-off\n"
        "9. CR + CAB approval  10. Deployment log  11. Post-deploy verify\n"
        "12. Security scan / pen test (if required)  13. RCA (if incident)",
    ),
    (
        "Key Metrics — Ops & Leadership Dashboard",
        "BRD first-pass acceptance rate          Target ≥ 60%\n"
        "BRD → triage decision time              Target < 5 days\n"
        "Deployments without linked BRD          Target 0\n"
        "UAT sign-off before production          Target 100%\n"
        "Emergency prod changes / month          Target ≤ 2\n"
        "Open IT-Security prod exceptions        Target 0",
    ),
    (
        "Escalation Paths",
        "Skip BRD pressure       → Ops → IT-Governance policy → IT Director\n"
        "GRC hold + deploy rush  → Document hold → GRC Manager → no deploy\n"
        "Security review open    → Ops blocks CR → IT-Security lead\n"
        "Scope dispute           → BA → IT Product → Sponsor\n"
        "CAB rejection           → IT-Governance → IT Product → Sponsor\n"
        "Material incident       → IT-Security IR → CIO → EXCO / BOD",
    ),
    (
        "Summary",
        "DEMAND: Business defines in BRD — IT does not guess rules\n"
        "DELIVER: Scrum builds from FSD with traceability to BRD\n"
        "ASSURE: IT-Governance + IT-Security + GRC run in parallel\n"
        "SHIP: CAB + UAT — then hypercare → BAU support\n\n"
        "Docs: docs/12-it-operations-stakeholder-framework.md\n"
        "      docs/11-operations-manager-checklist.md",
    ),
    (
        "Next Steps",
        "• Publish framework on Confluence (FECBRD space)\n"
        "• Configure Jira: Epic = BRD, Feature = FSD, link fields mandatory\n"
        "• Train dev leads on sprint planning gate checklist\n"
        "• Align CAB charter with IT-Governance + IT-Security\n"
        "• Weekly Ops metrics to IT leadership",
    ),
]


def _add_numbered_steps_slide(prs, title, steps, subtitle=None):
    """Two-column numbered step tiles for guides, restyled."""
    s, y_start = _frame(prs, title)
    if subtitle:
        _text(s, 0.7, y_start, 12.0, 0.35, [(subtitle, 12.5, False, GREY)])
        y_start += 0.5

    half = (len(steps) + 1) // 2
    col_w = 5.9
    left_x, right_x = 0.55, 6.85
    region_h = 6.85 - y_start
    rh = min(0.92, (region_h - 0.12 * (half - 1)) / half)
    for i, (num, text) in enumerate(steps):
        col = 0 if i < half else 1
        row = i if i < half else i - half
        x = left_x if col == 0 else right_x
        yy = y_start + row * (rh + 0.12)
        acc = ACCENTS[i % len(ACCENTS)]
        _rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, yy, col_w, rh, fill=CARD_FILL, line=BORDER, line_w=1)
        circle = _rect(s, MSO_SHAPE.OVAL, x + 0.18, yy + (rh - 0.46) / 2, 0.46, 0.46, fill=acc)
        _set_shape_text(circle, [(str(num), 13, True, CREAM)])
        _text(s, x + 0.8, yy, col_w - 1.0, rh, [(text, 11.5, False, INK)], anchor=MSO_ANCHOR.MIDDLE)

    return s


BUSINESS_USER_BRD_SLIDES = [
    (
        "How to Write a BRD at Finance",
        "A practical guide for business users\n"
        "Hướng dẫn viết BRD cho nhân viên nghiệp vụ\n\n"
        "Your BRD is the contract for what IT will deliver.",
    ),
    (
        "What You Will Learn",
        "• What a BRD is — and what it is NOT\n"
        "• The 8 steps to write a good BRD\n"
        "• Each template section (A → M) with examples\n"
        "• How to pass the quality gate (score ≥ 80%)\n"
        "• Where and how to submit",
    ),
    (
        "Why Your BRD Matters",
        "• IT cannot build what you have not defined\n"
        "• Poor BRD = delay, rework, UAT failure, compliance risk\n"
        "• Good BRD = faster triage, predictable delivery\n\n"
        "You own the business rules. IT owns the build.",
    ),
    (
        "What Is a BRD?",
        "A BRD describes:\n"
        "  ✓ Business problem and impact\n"
        "  ✓ Who is affected and business rules\n"
        "  ✓ Scope, data, compliance, acceptance criteria\n\n"
        "A BRD is NOT:\n"
        "  ✗ Technical design (API, database, Kafka)\n"
        "  ✗ Vendor selection or project plan\n"
        "  ✗ An email saying \"please do it ASAP\"",
    ),
    (
        "BRD vs Other Documents",
        "YOU write → BRD (business need)\n"
        "BA + IT write → FSD / FRD (functional spec)\n"
        "IT writes → Technical design & code\n"
        "YOU sign → UAT (did it meet your criteria?)\n\n"
        "No BRD → no IT build. No UAT sign-off → no go-live.",
    ),
    (
        "Your Journey After Submit",
        "1. You draft BRD + get Sponsor (Director+) to sign\n"
        "2. BA quality review (2 days) — score must be ≥ 80%\n"
        "3. Compliance routing if needed (Risk, Legal, IT-Security)\n"
        "4. IT triage & prioritization (5 days)\n"
        "5. BA creates FSD → Dev builds → You run UAT\n"
        "6. Go-live after your UAT sign-off\n\n"
        "Tip: Strong acceptance criteria in your BRD become your UAT tests.",
    ),
    (
        "The Golden Rule",
        "Write the PROBLEM and BUSINESS RULES,\n"
        "not the SYSTEM DESIGN.\n\n"
        "Viết VẤN ĐỀ và QUY TẮC NGHIỆP VỤ,\n"
        "không viết THIẾT KẾ HỆ THỐNG.",
    ),
    (
        "Good vs Bad Request",
        "BAD:\n"
        "  \"Build Kafka API to Finacle for real-time sync\"\n"
        "  \"Need new screen ASAP\"\n\n"
        "GOOD:\n"
        "  \"POS staff cannot see approved loan status during customer visit,\n"
        "   causing 18% drop-off on pre-approved applications\"\n\n"
        "BAD = IT solution.  GOOD = business problem + impact.",
    ),
    (
        "Problem Statement Formula",
        "[Role] cannot [task] because [constraint], causing [impact].\n\n"
        "[Vai trò] không thể [làm gì] vì [rào cản], gây [ảnh hưởng].\n\n"
        "Example / Ví dụ:\n"
        "Collections agent cannot see payment promise history on one screen,\n"
        "causing 12 min average call time and repeat contacts.",
    ),
    (
        "Section A — Request Information",
        "Fill in first:\n"
        "• BRD title (clear, specific)\n"
        "• Business unit (Sales/POS, Digital, Collections…)\n"
        "• Your name and email\n"
        "• Business Sponsor — must be Director grade or above\n"
        "• Target go-live date (realistic: ≥ 15 business days from submit)\n"
        "• Priority: Regulatory | Revenue | CX | Efficiency | Risk\n\n"
        "Common mistake: no Sponsor → BRD returned immediately.",
    ),
    (
        "Section B — Executive Summary",
        "Answer in plain language:\n"
        "• What is the business problem?\n"
        "• Who is impacted? (role, count, location)\n"
        "• What happens if we do nothing?\n"
        "• Desired outcome in one sentence\n\n"
        "Use the problem formula. Min 50 characters for problem statement.",
    ),
    (
        "Section C — Objectives & KPIs",
        "Every objective needs THREE parts:\n"
        "  Baseline (today) → Target (goal) → How measured\n\n"
        "WEAK: \"Improve customer experience\"\n"
        "STRONG: \"Reduce POS disbursement TAT from 24h to 4h;\n"
        "         measured weekly from LOS report R-1042\"",
    ),
    (
        "Section D — Current State",
        "Describe TODAY before IT changes anything:\n"
        "• Step-by-step current process\n"
        "• Which Finance systems are involved\n"
        "   (FE ONLINE 2.0, Finacle, POS/LOS, Collections…)\n"
        "• How many users and transactions per month\n\n"
        "IT uses this to understand where you are now.",
    ),
    (
        "Section E & F — To-Be & Scope",
        "Section E — Future process (business view only):\n"
        "  What should users be able to do differently?\n"
        "  Do NOT name APIs, databases, or vendors.\n\n"
        "Section F — Scope (write BOTH):\n"
        "  In scope: what this change delivers\n"
        "  Out of scope: what is NOT included (prevents scope creep)",
    ),
    (
        "Section G & H — Users & Business Rules",
        "Section G — Who uses it:\n"
        "  Role | Department | # users | Location (POS/HQ/field)\n\n"
        "Section H — Rules IT must NOT guess:\n"
        "  • Who approves? What authority level?\n"
        "  • Cut-off times (e.g. disbursement before 17:00)\n"
        "  • Exceptions and edge cases\n"
        "  • Customer communication rules",
    ),
    (
        "Section I — Data & Security",
        "Classify your data honestly:\n"
        "  Public | Internal | Confidential | Restricted\n\n"
        "Check data types: PII, CIC, payment, financial\n"
        "Remote access: None | Managed device | VDI | Exception\n\n"
        "Restricted data or exceptions → IT-Security review (adds time).\n"
        "Never request: disable MFA, personal Gmail, BYOD for customer data.",
    ),
    (
        "Section N — Compliance Screening",
        "Answer Yes / No / N/A honestly for all 7 questions:\n"
        "• Loan origination / approval / disbursement change?\n"
        "• Interest, fees, or contract terms?\n"
        "• Customer SMS / email / push notifications?\n"
        "• Third-party data sharing?\n"
        "• Collections or legal process?\n"
        "• Audit trail required?\n"
        "• Customer data on POS or field devices?\n\n"
        "\"Yes\" routes to Risk, Legal, or IT-Security — plan extra time.",
    ),
    (
        "Section M — Acceptance Criteria (UAT)",
        "Minimum 5 testable criteria using Given / When / Then:\n\n"
        "Given a pre-approved POS loan application\n"
        "When the agent opens the customer record during visit\n"
        "Then approved amount and status display within 3 seconds\n\n"
        "These become YOUR UAT tests. IT does not accept for you.",
    ),
    (
        "Section L & K — Risks & Rollout",
        "Section L — What goes wrong if built incorrectly?\n"
        "  Operational, customer, financial, compliance risks\n\n"
        "Section K — Rollout:\n"
        "  Pilot vs phased vs big bang\n"
        "  Training needed? Vietnamese / English?\n"
        "  Business continuity fallback if system unavailable",
    ),
    (
        "Section O — Sponsor Approval",
        "Before you submit:\n"
        "☐ Business Sponsor (Director+) has reviewed the full BRD\n"
        "☐ Sponsor signature or approval email attached\n"
        "☐ Sponsor cannot be the only approver for Critical changes\n\n"
        "Sponsor accountability = business owns the decision.",
    ),
    (
        "Top 10 Mistakes to Avoid",
        "1. Writing a solution instead of a need\n"
        "2. No KPI baseline or target\n"
        "3. Missing business rules (IT would have to guess)\n"
        "4. No out-of-scope section\n"
        "5. Multiple unrelated projects in one BRD\n"
        "6. Fewer than 5 acceptance criteria\n"
        "7. No Sponsor sign-off\n"
        "8. Ignoring POS / field training impact\n"
        "9. Wrong or missing data classification\n"
        "10. Copy-paste from vendor marketing material",
    ),
    (
        "Self-Check Before Submit",
        "☐ Can IT build without guessing my business rules?\n"
        "☐ Can UAT verify success using my acceptance criteria?\n"
        "☐ Can Risk/Audit understand control impact?\n"
        "☐ Sponsor (Director+) signed?\n"
        "☐ In-scope AND out-of-scope written?\n"
        "☐ Compliance Section N complete?\n"
        "☐ Quality preview ≥ 80% (use BRD Intake App)",
    ),
    (
        "Quality Gate — What BA Checks",
        "Pass requirements:\n"
        "• All mandatory sections complete\n"
        "• Quality score ≥ 80%\n"
        "• Sponsor sign-off on file\n\n"
        "Scoring tip: Business rules (Section H) weighted 20%.\n"
        "Score 0 on business rules → BRD usually fails.\n\n"
        "If returned: BA lists specific gaps — fix and resubmit within 10 days.",
    ),
    (
        "Compliance Quick Reference",
        "Loan approval / disbursement → Risk + Credit Policy\n"
        "Interest / fees / contract → Legal + Finance\n"
        "Customer SMS / email → Legal (approved template)\n"
        "Third-party data → Vendor Risk + Legal\n"
        "Collections / legal process → Compliance\n"
        "Data on POS / home device → IT-Security",
    ),
    (
        "Tools & Resources",
        "BRD Intake App (wizard + score preview):\n"
        "  app/ — run local server, export Markdown\n\n"
        "Templates & samples:\n"
        "  docs/01-brd-template-en.md\n"
        "  docs/02-brd-template-bilingual-vi-en.md\n"
        "  examples/04a-brd-pos-lending.md (gold standard)\n\n"
        "Submit via: ServiceNow / Jira BRD catalog (not email)",
    ),
    (
        "What NOT to Put in Your BRD",
        "✗ API names, Kafka, database, microservices\n"
        "✗ \"ASAP\" without business reason\n"
        "✗ Request to disable MFA or DLP\n"
        "✗ Personal device / Gmail / Dropbox for customer data\n"
        "✗ IT headcount or sprint planning\n\n"
        "If you need access or password reset → Service Desk, not BRD.",
    ),
    (
        "Summary — 8 Steps",
        "See visual on previous slide. Remember:\n"
        "1 Problem  2 Quantify  3 Who uses  4 Business rules\n"
        "5 Scope in/out  6 Data & compliance  7 Acceptance criteria  8 Sponsor sign\n\n"
        "Badge: BRD Ready — Finance (after training + 1 BRD ≥ 80%)",
    ),
    (
        "Next Steps",
        "1. Open BRD template or Intake App\n"
        "2. Draft using this guide — start with problem, not system\n"
        "3. Self-check and aim for ≥ 80% quality score\n"
        "4. Get Sponsor sign-off\n"
        "5. Submit via ServiceNow / Jira BRD form\n"
        "6. Attend BA office hours if you need coaching\n\n"
        "Questions: BA team | Confluence: FECBRD space",
    ),
]


def generate_business_user_brd_pptx():
    global FOOTER_TEXT, DECK_KICKER
    FOOTER_TEXT = "How to Write a BRD  ·  Business User Guide"
    DECK_KICKER = "Finance · Business User Guide"
    prs = Presentation()
    prs.slide_width = PInches(13.333)
    prs.slide_height = PInches(7.5)

    _add_content_slide(prs, *BUSINESS_USER_BRD_SLIDES[0])
    _add_content_slide(prs, *BUSINESS_USER_BRD_SLIDES[1])
    _add_content_slide(prs, *BUSINESS_USER_BRD_SLIDES[2])

    _add_pipeline_slide(
        prs,
        "Your BRD in the IT Delivery Pipeline",
        ["You write\nBRD", "BA review\n≥80%", "IT\ntriage", "FSD\nbuilt", "You test\nUAT", "Go-\nlive"],
        subtitle="Business users own BRD and UAT — IT owns build and deploy",
    )

    for item in BUSINESS_USER_BRD_SLIDES[3:7]:
        _add_content_slide(prs, *item)

    _add_numbered_steps_slide(
        prs,
        "8 Steps to a Good BRD",
        [
            (1, "Start with the problem — not \"we need system X\""),
            (2, "Quantify impact — %, hours, customers, VND"),
            (3, "Define who uses it — role, count, location"),
            (4, "Document business rules — approval, cut-off, exceptions"),
            (5, "Write in-scope AND out-of-scope"),
            (6, "Complete data & compliance honestly"),
            (7, "Add acceptance criteria — Given/When/Then (min 5)"),
            (8, "Get Sponsor (Director+) sign-off before IT intake"),
        ],
        subtitle="8 bước viết BRD tốt | Same steps in Vietnamese cheat sheet",
    )

    for item in BUSINESS_USER_BRD_SLIDES[7:]:
        if item[0] == "Summary — 8 Steps":
            continue  # covered by numbered steps slide
        _add_content_slide(prs, *item)

    out = OUTPUT / "Finance-Business-User-BRD-Guide-Slides.pptx"
    prs.save(out)
    print(f"Created {out}")
    return out


IT_OPS_SLIDES = [
    (
        "Finance IT Operations Guide",
        "Delivery control · Support tickets · Incidents · Daily monitoring\n"
        "For Ops Manager, Service Desk, on-call & release coordinators",
    ),
    (
        "Learning Objectives",
        "• Control delivery — enforce gates without doing business work\n"
        "• Triage support tickets to the correct bucket (IT vs Business)\n"
        "• Run BOD/EOD rituals and monitor systems throughout the day\n"
        "• Handle incidents with business liaison — IT restores systems first\n"
        "• Know when to redirect, escalate, and schedule post-mortems",
    ),
    (
        "Scope Boundary — IT vs Business",
        "IT OPERATIONS OWNS:\n"
        "  System availability · Incidents & recovery · Deploy / CAB / CR\n"
        "  Access / infra / env · Execute import AFTER business approval\n\n"
        "BUSINESS OWNS (not IT):\n"
        "  Case / loan / collection updates · Business rules & decisions\n"
        "  Data content & maker-checker · UAT sign-off · BRD & Sponsor\n\n"
        "Service Desk script:\n"
        "  \"IT fixes systems and access. Business owns case updates and rules.\n"
        "   I'll route you to the right channel.\"",
    ),
    (
        "Golden Rules for Ops",
        "1. IT executes; business defines — no build without accepted BRD (≥ 80%)\n"
        "2. Email / chat are NOT intake for projects or rule changes\n"
        "3. Data import is a business act — BRD + approved file + maker-checker\n"
        "4. UAT sign-off is business-owned — IT does not accept for them\n"
        "5. No production data fix without approved business instruction + ticket\n"
        "6. Document and redirect wrong-bucket requests — do not silently help",
    ),
    (
        "Shift Schedule — Ops & Service Desk",
        "WEEKDAY (Mon–Fri):\n"
        "  NIGHT SHIFT      20:00 – 24:00   Monitor · batch · on-call · handover to BOD\n"
        "  BOD WINDOW       04:00 – 08:00   Health · releases · ready-state prep\n"
        "  BUSINESS ALLOC   07:30           Business users allocate apps — systems green\n"
        "  DAY SHIFT        08:30 – 17:30   Catalog SR · incidents · delivery control\n"
        "  EOD              17:30           Close day · handover to night (20:00)\n\n"
        "Rule: Ops completes BOD before 07:30 so business can allocate on time.",
    ),
    (
        "Ops Roster — Who Handles What",
        "ROSTER BY WORK TYPE (not by who asks loudest):\n\n"
        "IT OPS / SERVICE DESK (catalog SR + incidents):\n"
        "  Access · system down/slow · deploy/CR · env · approved import execution\n\n"
        "BUSINESS OPS (self-service — IT redirects, does NOT execute):\n"
        "  Case updates · loan/collection decisions · allocation disputes\n"
        "  Rule changes · bulk data · feature requests · UAT sign-off\n\n"
        "CASE HANDLING:\n"
        "  IT provides logs & system restore only\n"
        "  Business unit owns case outcome — delegate back same day\n\n"
        "Publish weekly roster: name | shift | escalation backup",
    ),
    (
        "Weekend Coverage — Limited Hours",
        "Weekends: NO full day shift. BOD + EOD + short day window only.\n\n"
        "STANDARD WEEKEND (no trigger):\n"
        "  BOD (04:00–08:00) · on-call monitor only · EOD handover\n\n"
        "EXTENDED DAY SHIFT (max 4 hours) — only if:\n"
        "  ☐ Batch/data feed late — business cannot allocate\n"
        "  ☐ P1/P2 incident affecting allocation or catalog access\n"
        "  ☐ Critical SR in IT catalog (access/system) — not feature delivery\n\n"
        "OUT OF SCOPE weekends: new features, BRD work, case updates by IT\n"
        "  → Redirect to business self-service · log wrong_bucket\n\n"
        "Ops Manager approves extended weekend roster before shift starts.",
    ),
    (
        "IT Catalog SR — Support Only, Not Delivery",
        "IT Service Catalog accepts ONLY operational support:\n\n"
        "IN CATALOG (IT resolves):\n"
        "  Password / VPN / access · account lock · system error\n"
        "  Performance · integration down · approved import run\n\n"
        "NOT IN CATALOG (delegate to business — close IT ticket):\n"
        "  Update case / status / allocation · change business rules\n"
        "  New screen / report / feature · \"build this for us\"\n"
        "  Explain why loan rejected · fix customer outcome\n\n"
        "Script: \"This is business-owned work. IT catalog is for systems\n"
        "  and access. Please use your business process / BRD intake.\"\n\n"
        "Feature delivery → BRD → triage — never Service Desk queue.",
    ),
    (
        "Delivery Control — Your Role",
        "Ops Manager controls the PIPELINE, not business decisions:\n\n"
        "☐ Enforce one front door (ITSM catalog)\n"
        "☐ Block IT execution on business-bucket work\n"
        "☐ Verify BRD linked before FSD / sprint / deploy\n"
        "☐ Run release readiness + CAB evidence check\n"
        "☐ Own hypercare (T+1 to T+14) after go-live\n"
        "☐ Sample tickets daily for bucket correctness\n\n"
        "You say NO to deploy — you do not update customer cases.",
    ),
    (
        "Release Readiness — Go / No-Go",
        "Before ANY production deploy, verify:\n\n"
        "☐ BRD accepted (score ≥ 80%) and linked to CR\n"
        "☐ FSD complete; in-scope stories closed\n"
        "☐ SIT exit criteria met — test summary attached\n"
        "☐ UAT signed by Business Sponsor (not IT)\n"
        "☐ CAB approved (IT-Governance)\n"
        "☐ IT-Security sign-off (if flagged)\n"
        "☐ Rollback plan tested · Comms & training done\n"
        "☐ Hypercare roster assigned\n\n"
        "No-go = document reason; escalate to IT Product + Sponsor",
    ),
    (
        "Hard Gates — Non-Negotiable",
        "1. No FSD / sprint without accepted BRD (≥ 80%)\n"
        "2. No production deploy without UAT + CAB\n"
        "3. No security bypass without IT-Security exception register\n"
        "4. Emergency change → retro-BRD within 2 business days\n"
        "5. Deployments without linked BRD = 0 (monthly KPI)\n"
        "6. Business data fixes by IT without approval = 0",
    ),
    (
        "Bucket Allocation Matrix",
        "PASSWORD / VPN / ACCESS          → IT bucket · Service Desk\n"
        "SYSTEM DOWN / ERROR / SLOW         → IT bucket · Incident\n"
        "CHANGE APPROVAL RULE / ADD FIELD   → Business bucket · BRD intake\n"
        "UPDATE 500 CASE STATUSES           → Business bucket · import approval\n"
        "BUILD FEATURE / INTEGRATION        → Business bucket · BRD → triage\n"
        "DEPLOY RELEASE                     → IT bucket · CR after UAT + CAB\n"
        "IMPORT PARTNER FILE                → Business approval → IT executes\n"
        "WHY WAS LOAN REJECTED?             → Business ops — IT provides logs only",
    ),
    (
        "30-Second Triage Card",
        "Ask in order — stop when answered:\n\n"
        "1. Is the SYSTEM broken, slow, or access blocked?\n"
        "   YES → IT ticket (Incident or Service Desk)\n\n"
        "2. Is it CHANGE RULE / CASE / IMPORT / REPORT DEFINITION?\n"
        "   YES → Business bucket → redirect to BRD\n\n"
        "3. Is it BUILD / ENHANCE / INTEGRATE?\n"
        "   YES → Accepted BRD linked? NO → redirect to BRD first\n\n"
        "4. Still unsure? → BA 30-min triage call — DO NOT assign to Dev",
    ),
    (
        "Support Ticket Runbook",
        "IT catalog SR only — NOT feature delivery or case work:\n\n"
        "1. ACKNOWLEDGE (≤ 15 min P3) — ITSM ticket from catalog\n"
        "2. CLASSIFY — IT catalog / Business self-service / Incident (≤ 30 min)\n"
        "3. If business-owned (case, rule, feature) → redirect & close IT queue\n"
        "4. If IT catalog → assign group + priority (≤ 1 hr)\n"
        "5. RESOLVE per SLA — systems & access only\n"
        "6. CLOSE — user confirms; KB article if repeat\n"
        "7. EOD — wrong-bucket count; weekend triggers logged",
    ),
    (
        "Redirect Template — Delegate to Business",
        "Send when request is NOT in IT catalog (case, rule, feature):\n\n"
        "\"This request is business-owned — IT does not update cases,\n"
        "  allocation, or rules via Service Desk.\n\n"
        " IT catalog is for: access, system errors, approved imports.\n\n"
        " Your team must handle:\n"
        "  • Case / allocation updates → business ops process\n"
        "  • Rule or feature changes → BRD intake + Sponsor sign-off\n"
        "  • Outcome disputes (loan rejected, etc.) → Credit / business ops\n\n"
        " IT can provide audit logs only. Please use your business channel.\"",
    ),
    (
        "Incident vs SR vs Business Work",
        "INCIDENT (IT): FE ONLINE down; POS cannot submit; batch stuck\n"
        "SERVICE REQUEST (IT): Password reset; new user access\n"
        "BUSINESS WORK (not IT): Fix wrong disbursement; update case status\n"
        "DEFECT (IT after release): Bug in new feature → link to release\n\n"
        "User says \"incident\" but means business activity?\n"
        "  \"Update 200 cases to closed\" → Business import process\n"
        "  \"Customer angry — fix their loan\" → Business ops; IT = audit log only",
    ),
    (
        "Incident Priority — P1 to P4",
        "P1 CRITICAL — major business stop (core banking down, mass POS fail)\n"
        "   Notify within 30 min: Ops, IT-Security, CIO, EXCO\n\n"
        "P2 MAJOR — significant degradation (login fail, collections batch stuck)\n"
        "   Notify within 1 hr: Ops, IT Product, Sponsor\n\n"
        "P3 MODERATE — workaround exists (single region slow)\n"
        "   Notify within 4 hr\n\n"
        "P4 MINOR — cosmetic / single user → next business day",
    ),
    (
        "Incident Response — Phases",
        "DETECT     Monitor + user report → log incident in ITSM\n"
        "TRIAGE     Classify P1–P4; open war room if P1/P2\n"
        "DIAGNOSE   Logs, metrics, recent CRs — IT only\n"
        "CONTAIN    Rollback, disable feature, scale — document decision\n"
        "RESOLVE    Fix deploy or data repair WITH business approval only\n"
        "RECOVER    Verify monitoring green; business confirms ops resumed\n"
        "CLOSE      Timeline in ITSM; post-mortem if P1/P2 or repeat",
    ),
    (
        "War Room Checklist — P1 / P2",
        "☐ Incident commander appointed (Ops Manager)\n"
        "☐ Timeline log — 5-minute updates\n"
        "☐ IT-Security notified if data / access impact\n"
        "☐ Business liaison on bridge (does NOT decide IT technical actions)\n"
        "☐ GRC notified if regulatory impact\n"
        "☐ Customer comms draft — Legal if customer-facing\n"
        "☐ Rollback decision documented\n"
        "☐ Post-mortem scheduled within 5 business days",
    ),
    (
        "System Monitoring — All Day",
        "Weekday: day shift (08:30–17:30) + night (20:00–24:00) on-call:\n\n"
        "☐ Production health dashboard (green / amber / red)\n"
        "☐ Alert queue — no unacknowledged critical alerts > 15 min\n"
        "☐ Batch / EOD jobs — lending, collections, GL cut-offs\n"
        "☐ Integration queues — Finacle, FE ONLINE 2.0, POS\n"
        "☐ Hypercare releases — daily Sponsor check (T+1 to T+14)\n"
        "☐ Capacity — CPU, disk, connection pools on critical paths\n"
        "☐ Security — failed login spikes, WAF blocks (escalate IT-Security)\n\n"
        "Amber → investigate. Red → incident process.",
    ),
    (
        "EOD Runbook — Day Shift Close 17:30",
        "End of day shift (08:30–17:30) — handover to night (20:00–24:00):\n\n"
        "☐ All P1/P2 updated in ITSM (not stale)\n"
        "☐ Open SRs — none unassigned past SLA\n"
        "☐ Business-bucket redirects logged with reason\n"
        "☐ Deployments today — post-deploy verify complete?\n"
        "☐ Emergency changes — retro-BRD task opened?\n"
        "☐ Handover note for night / on-call\n"
        "☐ No informal prod changes via chat — audit sample\n"
        "☐ Update Ops dashboard metrics (daily slice)",
    ),
    (
        "BOM & EOM — Monthly Rituals",
        "BOM (1st business day, 60 min):\n"
        "  Last month KPIs · BRD volume vs BA capacity · release calendar\n"
        "  Incident trends · wrong-bucket training plan · CAB roster\n\n"
        "EOM (last business day, 60 min):\n"
        "  Change freeze for critical systems · evidence pack sample\n"
        "  Deploys without BRD = 0 · UAT before prod = 100%\n"
        "  Incident summary to leadership · lessons learned top 3\n"
        "  Wrong-bucket report to BU Heads",
    ),
    (
        "Post-Mortem & Lessons Learned",
        "Mandatory when: P1/P2 · failed deploy · repeat P3 (3× in 90 days)\n"
        "  · emergency change · wrong-bucket caused customer impact\n\n"
        "Blameless focus: timeline of facts → root cause → SMART actions\n"
        "Separate: IT systems failure vs business process vs handoff confusion\n\n"
        "Distribute lessons to: Service Desk KB · BU newsletter · Dev retro\n"
        "  · EOM leadership pack · quarterly audit",
    ),
    (
        "Ops Guardrails — Posture Card",
        "Post at Service Desk and war room:\n\n"
        "1. Systems yes — business decisions no\n"
        "2. No prod data fix without approved instruction + ticket\n"
        "3. No rule change without BRD + FSD + deploy path\n"
        "4. No import without maker-checker\n"
        "5. Incidents — restore service first; data correction follows business\n"
        "6. Email is not intake for changes\n"
        "7. Document and redirect wrong-bucket — don't do business work silently\n"
        "8. P1/P2 — post-mortem always",
    ),
    (
        "Key KPIs — Ops Dashboard",
        "Wrong-bucket SR rate              → trend down (EOD / EOM)\n"
        "SR classified within 1 hr       → target ≥ 95%\n"
        "P1 MTTR                         → per policy (EOM)\n"
        "Post-mortems on time (P1/P2)    → target 100%\n"
        "Deploys without BRD (month)     → target 0\n"
        "Repeat incidents (same RCA)     → target 0 per quarter\n"
        "IT data fixes without approval  → target 0",
    ),
    (
        "Summary",
        "WEEKDAY:  Night 20:00–24:00 · BOD 04:00–08:00 · Biz 07:30 · Day 08:30–17:30 · EOD 17:30\n"
        "WEEKEND:  BOD + EOD only · day shift max 4h if data late or allocation/incident\n"
        "ROSTER:   IT = catalog SR + systems · Business = cases, rules, features (self-service)\n"
        "TICKETS:  Catalog = support only — NOT feature delivery · redirect wrong-bucket\n"
        "MONTHLY:  BOM capacity · EOM audit · lessons learned\n\n"
        "IT fixes systems. Business owns cases, allocation, rules, and data decisions.",
    ),
    (
        "Resources & Next Steps",
        "Runbook:     docs/14-it-operations-runbook.md\n"
        "Checklist:   docs/11-operations-manager-checklist.md\n"
        "Framework:   docs/12-it-operations-stakeholder-framework.md\n"
        "SO guide:    docs/13-service-owner-delivery-control-guide.md\n"
        "BRD App:     app/ (intake wizard for business redirects)\n\n"
        "Actions:\n"
        "• Print 30-second triage card at Service Desk\n"
        "• Run BOD/EOD logs in ITSM or Confluence daily\n"
        "• Weekly wrong-bucket report to IT leadership",
    ),
]


def generate_it_operations_pptx():
    global FOOTER_TEXT, DECK_KICKER
    FOOTER_TEXT = "IT Operations Guide  ·  Internal use only"
    DECK_KICKER = "Finance · IT Operations"
    prs = Presentation()
    prs.slide_width = PInches(13.333)
    prs.slide_height = PInches(7.5)

    _add_content_slide(prs, *IT_OPS_SLIDES[0])
    _add_content_slide(prs, *IT_OPS_SLIDES[1])
    _add_content_slide(prs, *IT_OPS_SLIDES[2])

    _add_pipeline_slide(
        prs,
        "Daily Ops Rhythm",
        [
            "BOD\n04-08",
            "Biz alloc\n07:30",
            "Day shift\n08:30-17:30",
            "Monitor\n& SR",
            "EOD\n17:30",
        ],
        subtitle="Night 20:00–24:00 · Weekend: BOD/EOD + max 4h day if triggered",
    )

    for item in IT_OPS_SLIDES[3:8]:
        _add_content_slide(prs, *item)

    _add_pipeline_slide(
        prs,
        "Delivery Control Pipeline",
        [
            "Receive\nReq",
            "BRD\nGate",
            "Triage\n& FSD",
            "Build\n& SIT",
            "UAT\nsign",
            "CAB\nShip",
            "Hyper-\ncare",
        ],
        subtitle="Ops enforces gates — does not skip UAT or CAB under pressure",
    )

    for item in IT_OPS_SLIDES[8:12]:
        _add_content_slide(prs, *item)

    _add_pipeline_slide(
        prs,
        "Support Ticket Handling Flow",
        ["Catalog\nITSM", "IT or\nBusiness?", "IT SR\nresolve", "Business\nself-svc", "Close /\nredirect"],
        subtitle="Catalog = support only · case/rule/feature → delegate to business",
    )

    # VISUAL: decision flowchart for ticket triage
    _add_flowchart_slide(
        prs,
        "Flowchart — Ticket Triage Decision",
        [
            {"kind": "start", "text": "Request arrives (ITSM catalog)"},
            {"kind": "decision", "text": "System broken / slow / access blocked?",
             "branch": ("Yes", "IT bucket → Incident or Service Desk (resolve per SLA)")},
            {"kind": "decision", "text": "Change rule / case / allocation / import?",
             "branch": ("Yes", "Business bucket → redirect to BRD / self-service")},
            {"kind": "decision", "text": "Build / enhance / new feature?",
             "branch": ("Yes", "BRD accepted ≥80%? No → redirect; Yes → IT triage")},
            {"kind": "end", "text": "Unsure → BA 30-min triage (no Dev assign)"},
        ],
        subtitle="Stop at the first Yes — IT executes systems only, never business work",
    )

    for item in IT_OPS_SLIDES[12:16]:
        _add_content_slide(prs, *item)

    _add_pipeline_slide(
        prs,
        "Incident Response Flow",
        ["Detect\n& log", "Triage\nP1-P4", "Diagnose\n& contain", "Resolve\n& verify", "Post-\nmortem"],
        subtitle="Restore service first — business data correction follows approved process",
    )

    for item in IT_OPS_SLIDES[16:19]:
        _add_content_slide(prs, *item)

    _add_numbered_steps_slide(
        prs,
        "BOD Runbook — Window 04:00 to 08:00",
        [
            (1, "04:00 — Open BOD; handover from night shift (20:00–24:00 prior day)"),
            (2, "Review overnight P1/P2 — none open or unowned"),
            (3, "Confirm production health + overnight batch jobs complete"),
            (4, "List releases today — CR approved? UAT linked?"),
            (5, "CAB outcomes — failed deploy? Hypercare check"),
            (6, "Wrong-bucket SR count; BRD queue > 5 days? Escalate"),
            (7, "07:30 — Confirm apps green for business user allocation"),
            (8, "Weekend? Approve extended day roster only if data late / incident"),
            (9, "08:00 — Publish BOD log; assign day-shift escalation path"),
            (10, "08:30 — Day shift (08:30–17:30) assumes Service Desk lead"),
        ],
        subtitle="Weekend: BOD + EOD; day max 4h only on trigger · Ops Manager approves roster",
    )

    # VISUAL: org chart — operations & support resourcing
    _add_org_chart_slide(
        prs,
        "Org Chart — Ops & Support Resourcing",
        "IT Ops Manager",
        [
            ("Service Desk Lead", ["L1 catalog SR", "Bucket triage", "Redirect / KB"]),
            ("Shift / On-call", ["Day 08:30–17:30", "Night 20:00–24:00", "Weekend on-call"]),
            ("Release & Incident", ["CAB / change", "Incident commander", "Hypercare / Dev liaison"]),
        ],
        subtitle="Roster by work type — IT handles systems; business work is delegated back",
    )

    for item in IT_OPS_SLIDES[19:]:
        _add_content_slide(prs, *item)

    out = OUTPUT / "Finance-IT-Operations-Guide-Slides.pptx"
    prs.save(out)
    print(f"Created {out}")
    return out


def generate_it_delivery_framework_pptx():
    global FOOTER_TEXT, DECK_KICKER
    FOOTER_TEXT = "IT Delivery Framework  ·  Internal use only"
    DECK_KICKER = "Finance · IT Delivery Framework"
    prs = Presentation()
    prs.slide_width = PInches(13.333)
    prs.slide_height = PInches(7.5)

    # Slide 1 — title (content)
    _add_content_slide(prs, *IT_DELIVERY_SLIDES[0])
    # Slide 2 — objectives
    _add_content_slide(prs, *IT_DELIVERY_SLIDES[1])

    # Slide 3 — VISUAL: end-to-end pipeline
    _add_pipeline_slide(
        prs,
        "End-to-End Delivery Pipeline",
        [
            "Receive\nReq",
            "BRD\nGate",
            "Triage\n& FSD",
            "Sprint\nBuild",
            "SIT\nTest",
            "UAT",
            "Ship\n(CAB)",
            "Ops\nSupport",
        ],
        subtitle="Governed pipeline with Agile/Scrum inside Build → SIT phases",
    )

    # VISUAL: cross-functional swimlane (who does what across phases)
    _add_swimlane_slide(
        prs,
        "Swimlane — Cross-Functional Delivery",
        ["Intake", "Define", "Design", "Build", "Test", "Release", "Operate"],
        [
            ("Business / Sponsor", ["Raise need", "Write BRD\n+ sign", None, None, "UAT\nsign-off", "Go-live\napprove", "Confirm\nservice"]),
            ("BA", ["Classify", "Score BRD\n≥80%", "Author\nFSD", "Refine\nstories", "Support\nUAT", None, None]),
            ("IT Product (PO)", [None, "Accept to\ntriage", "Prioritize\nbacklog", "Accept\nsprint", None, "Release\ndecision", None]),
            ("Scrum / Dev", [None, None, "Design\n+ tasks", "Build +\nunit test", "Fix SIT\ndefects", "Deploy", "Hotfix /\ndefect"]),
            ("QA", [None, None, "Test\nplan", "Prep\ncases", "Run SIT\n+ regress", None, None]),
            ("IT Ops / Assurance", ["Route\nticket", "Compliance\nrouting", "ARB /\nSecurity", None, None, "CAB +\nverify", "Monitor /\nhypercare"]),
        ],
        subtitle="No hand-off skips a lane — every phase has a named owner",
    )

    # Slides 4–10 content (documents through FSD)
    for item in IT_DELIVERY_SLIDES[2:9]:
        _add_content_slide(prs, *item)

    # Slide 11 — VISUAL: Agile overlay pipeline
    _add_pipeline_slide(
        prs,
        "Agile/Scrum Inside the Pipeline",
        [
            "Product\nBacklog",
            "Sprint\nPlanning",
            "Daily\nBuild",
            "SIT in\nSprint",
            "Sprint\nReview",
            "UAT\nWindow",
            "Release\nCAB",
            "Retro",
        ],
        subtitle="FSD feeds Product Backlog; Ship requires governance gates green",
    )

    for item in IT_DELIVERY_SLIDES[9:14]:
        _add_content_slide(prs, *item)

    # Scrum ceremonies visual
    _add_scrum_cycle_slide(prs)

    for item in IT_DELIVERY_SLIDES[14:20]:
        _add_content_slide(prs, *item)

    # Dual-track diagram
    _add_dual_track_slide(prs)

    # VISUAL: org chart — development + operations resourcing
    _add_org_chart_slide(
        prs,
        "Org Chart — Development & Operations",
        "CIO / IT Director",
        [
            ("IT Product (PO)", ["BA Team", "Scrum / Dev Team", "QA / Test"]),
            ("IT Ops Manager", ["Service Desk", "On-call / Night shift", "Release / CAB coord"]),
            ("IT-Governance & Security", ["CAB / Change", "IT-Security", "GRC / Legal liaison"]),
        ],
        subtitle="Delivery (build) and Operations (run) report into IT leadership — assurance is independent",
    )

    for item in IT_DELIVERY_SLIDES[20:]:
        _add_content_slide(prs, *item)

    out = OUTPUT / "Finance-IT-Delivery-Framework-Slides.pptx"
    prs.save(out)
    print(f"Created {out}")
    return out


def generate_pptx():
    global FOOTER_TEXT, DECK_KICKER
    FOOTER_TEXT = "Finance BRD Training  ·  Internal use only"
    DECK_KICKER = "Finance BRD Training"
    prs = Presentation()
    prs.slide_width = PInches(13.333)
    prs.slide_height = PInches(7.5)

    for title, body in SLIDES:
        _add_content_slide(prs, title, body)

    out = OUTPUT / "Finance-BRD-Training-Slides.pptx"
    prs.save(out)
    print(f"Created {out}")
    return out


if __name__ == "__main__":
    generate_brd_template_docx()
    generate_cheat_sheet_docx()
    generate_pptx()
    generate_business_user_brd_pptx()
    generate_it_delivery_framework_pptx()
    generate_it_operations_pptx()
    print("Done.")
