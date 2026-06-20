#!/usr/bin/env python3
"""Generate AI Skills Learning deck — definition, sources, practice, exercise per skill."""

import sys
from pathlib import Path
from repo_paths import EXPORTS_LEARNING

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches as PInches

ROOT = Path(__file__).resolve().parent
sys.path.insert(0, str(ROOT))
from anthropic_theme import (  # noqa: E402
    SW, SH, CLAY, SKY, OLIVE, FIG, CACTUS, SLATE, CREAM, INK, LIGHT, BORDER, CARD_FILL,
    blank as _blank, rect as _rect, text as _text, header as _header, footer as _footer,
    title_slide, summary_slide,
)
from learning_loader import skills_for_slides  # noqa: E402

OUTPUT = EXPORTS_LEARNING
OUTPUT.mkdir(parents=True, exist_ok=True)
FOOTER = "AI Skills Lab  ·  learning_data.py  ·  See Learning-Master-Slides.pptx for weekly route"

SKILLS = skills_for_slides()


def _panel(slide, x, y, w, h, accent, label, body):
    _rect(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, fill=CARD_FILL, line=BORDER, line_w=1)
    _rect(slide, MSO_SHAPE.RECTANGLE, x, y, w, 0.38, fill=accent)
    _text(slide, x + 0.12, y + 0.05, w - 0.24, 0.3, [(label, 10, True, CREAM)])
    _text(slide, x + 0.14, y + 0.48, w - 0.28, h - 0.55, [(body, 9.5, False, INK)])


def s_how_to_use(prs, idx):
    s = _blank(prs)
    y = _header(s, "How to use this deck", kicker="Learn → do → prove")
    panels = [
        (CLAY, "1 · DEFINE", "Read the definition slide. Say it aloud in your own words."),
        (SKY, "2 · READ", "Complete ★ sources that day — book chapter or course lesson."),
        (OLIVE, "3 · PRACTICE", "Type code yourself. No copy-paste. Banking lens always."),
        (FIG, "4 · EXERCISE", "Finish the lab task. Commit to GitHub with a clear message."),
    ]
    cw, ch, gap = 5.85, 1.55, 0.4
    x0 = (SW - (cw * 2 + gap)) / 2
    for i, (acc, head, sub) in enumerate(panels):
        r, c = divmod(i, 2)
        _panel(s, x0 + c * (cw + gap), y + 0.35 + r * (ch + 0.28), cw, ch, acc, head, sub)
    _text(
        s, 0.55, y + 3.55, 12.2, 0.55,
        [("Only advance when checkpoint ✓ passes — not when video is finished.", 13, True, SLATE)],
        align=PP_ALIGN.CENTER,
    )
    _footer(s, idx, FOOTER)
    return s


def s_skill(prs, idx, skill):
    s = _blank(prs)
    acc = skill["accent"]
    y = _header(s, skill["name"], kicker=f"{skill['weeks']} · AI Skills Lab")
    margin = 0.55
    gap = 0.28
    pw = (SW - 2 * margin - gap) / 2
    ph = 2.05
    top = y + 0.25
    _panel(s, margin, top, pw, ph, acc, "DEFINITION", skill["definition"])
    _panel(s, margin + pw + gap, top, pw, ph, SKY, "READ (books & sources)", skill["sources"])
    _panel(s, margin, top + ph + gap, pw, ph, OLIVE, "PRACTICE", skill["practice"])
    _panel(s, margin + pw + gap, top + ph + gap, pw, ph, CLAY, "EXERCISE", skill["exercise"])
    bar_y = top + 2 * ph + gap + 0.2
    _rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, margin, bar_y, SW - 2 * margin, 0.72, fill=LIGHT, line=acc, line_w=1.5)
    _rect(s, MSO_SHAPE.RECTANGLE, margin, bar_y, 1.35, 0.72, fill=acc)
    _text(s, margin + 0.08, bar_y + 0.18, 1.2, 0.4, [("CHECK", 10, True, CREAM)], align=PP_ALIGN.CENTER)
    _text(s, margin + 1.5, bar_y + 0.12, SW - 2 * margin - 1.6, 0.5, [(skill["checkpoint"], 10, True, INK)])
    _footer(s, idx, FOOTER)
    return s


def s_index(prs, idx):
    s = _blank(prs)
    y = _header(s, f"{len(SKILLS)} skills — index", kicker="In learning order · see Learning-Master-Slides.pptx for weeks")
    cols = 3
    cw, ch, gx, gy = 3.85, 0.72, 0.28, 0.14
    x0 = (SW - (cw * cols + gx * (cols - 1))) / 2
    for i, sk in enumerate(SKILLS):
        if i >= 21:
            break
        r, c = divmod(i, cols)
        x = x0 + c * (cw + gx)
        yy = y + 0.3 + r * (ch + gy)
        acc = sk["accent"]
        _rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, yy, cw, ch, fill=CARD_FILL, line=BORDER)
        _rect(s, MSO_SHAPE.RECTANGLE, x, yy, 0.1, ch, fill=acc)
        _text(s, x + 0.2, yy + 0.08, 0.35, 0.35, [(str(i + 1), 11, True, acc)])
        _text(s, x + 0.55, yy + 0.06, cw - 0.65, 0.55, [(sk["name"], 9.5, True, INK)])
    _footer(s, idx, FOOTER)
    return s


def generate():
    prs = Presentation()
    prs.slide_width = PInches(SW)
    prs.slide_height = PInches(SH)

    title_slide(
        _blank(prs),
        "AI Skills Lab",
        "Learn Each Skill in Detail",
        [
            "Definition · Books & sources · Practice · Exercise · Checkpoint",
            "19 skills · banking domain · zero Python → job-ready",
            "Companion: ai-skills-catalog.md · lab/",
        ],
    )

    s_how_to_use(prs, 2)
    s_index(prs, 3)
    for i, sk in enumerate(SKILLS, start=4):
        s_skill(prs, i, sk)

    summary_slide(
        _blank(prs),
        "Summary",
        "The 4-step loop for every skill",
        [
            (CLAY, "Define", "Understand before coding — say it in plain English"),
            (SKY, "Read", "One ★ source per session — book or course"),
            (OLIVE, "Practice", "Type code · banking exercise · git commit"),
            (FIG, "Prove", "Pass checkpoint before next skill"),
        ],
        "Start today: Skill 1 → week01_brd_checklist.py · reading-path.md Step 1",
    )

    out = OUTPUT / "AI-Skills-Learning-Slides.pptx"
    prs.save(out)
    n = len(prs.slides._sldIdLst)
    print(f"Created {out} ({n} slides)")
    return out


if __name__ == "__main__":
    generate()
