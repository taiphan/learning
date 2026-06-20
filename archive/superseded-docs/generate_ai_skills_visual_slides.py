#!/usr/bin/env python3
"""Visual AI Skills deck — journey map, phases, loop, skill flow."""

import sys
from pathlib import Path
from repo_paths import EXPORTS_LEARNING

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches as PInches

ROOT = Path(__file__).resolve().parent
sys.path.insert(0, str(ROOT))
from anthropic_theme import (  # noqa: E402
    SW, SH, CLAY, SKY, OLIVE, FIG, CACTUS, SLATE, CREAM, INK, GREY, LIGHT, BORDER, CARD_FILL,
    blank as _blank, rect as _rect, text as _text, header as _header, footer as _footer,
    title_slide, summary_slide,
)
from generate_ai_skills_slides import SKILLS  # noqa: E402

OUTPUT = EXPORTS_LEARNING
OUTPUT.mkdir(parents=True, exist_ok=True)
FOOTER = "AI Skills Visual  ·  Workbook + Canvas"

PHASES = [
    (CLAY, "Foundation", "W1–12", "Python · Git · SQL · pandas · EDA", 1, 5),
    (SKY, "Classical ML", "W13–22", "Models · Metrics · SHAP", 6, 8),
    (FIG, "GenAI & Agents", "W25–38", "Embed · RAG · LangGraph · Eval", 9, 13),
    (CACTUS, "Production", "W39–44", "FastAPI · Docker · Guardrails · CI", 14, 17),
    (OLIVE, "Career", "W47–52", "Use cases · Portfolio · Apply", 18, 19),
]


def s_journey(prs, idx):
    s = _blank(prs)
    y = _header(s, "12-month skill journey", kicker="Visual map · 19 skills")
    bar_y = y + 0.5
    total_w = 11.5
    x0 = 0.55
    seg_w = total_w / len(PHASES)
    for i, (acc, name, when, skills, _, _) in enumerate(PHASES):
        x = x0 + i * seg_w
        _rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, x + 0.05, bar_y, seg_w - 0.1, 1.35, fill=acc)
        _text(s, x + 0.12, bar_y + 0.08, seg_w - 0.24, 0.35, [(when, 10, True, CREAM)])
        _text(s, x + 0.12, bar_y + 0.42, seg_w - 0.24, 0.45, [(name, 11, True, CREAM)], align=PP_ALIGN.CENTER)
        _text(s, x + 0.12, bar_y + 0.88, seg_w - 0.24, 0.4, [(skills, 8, False, CREAM)], align=PP_ALIGN.CENTER)
    _text(s, 0.55, bar_y + 1.55, 12.2, 0.5,
          [("Start → Skill 1 Python → pass checkpoint → next skill. Never skip exercises.", 12, False, GREY)],
          align=PP_ALIGN.CENTER)
    _footer(s, idx, FOOTER)
    return s


def s_loop(prs, idx):
    s = _blank(prs)
    y = _header(s, "The 4-step loop", kicker="Every skill · every session")
    steps = [
        (CLAY, "WHAT", "Define\nunderstand skill"),
        (SKY, "READ", "★ sources\nthat day"),
        (OLIVE, "PRACTICE", "Type code\nbanking lens"),
        (FIG, "EXERCISE", "Lab task\ngit commit"),
        (CACTUS, "CHECK", "Pass before\nnext skill"),
    ]
    cw, gap = 2.15, 0.22
    x0 = (SW - (cw * len(steps) + gap * (len(steps) - 1))) / 2
    cy = y + 1.0
    for i, (acc, head, sub) in enumerate(steps):
        x = x0 + i * (cw + gap)
        _rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, cy, cw, 1.65, fill=CARD_FILL, line=acc, line_w=2)
        _rect(s, MSO_SHAPE.RECTANGLE, x, cy, cw, 0.42, fill=acc)
        _text(s, x + 0.08, cy + 0.06, cw - 0.16, 0.32, [(head, 10, True, CREAM)], align=PP_ALIGN.CENTER)
        _text(s, x + 0.1, cy + 0.55, cw - 0.2, 0.95, [(sub, 10, False, INK)], align=PP_ALIGN.CENTER)
        if i < len(steps) - 1:
            ax = x + cw + 0.02
            _text(s, ax, cy + 0.65, gap + 0.05, 0.35, [("→", 14, True, GREY)], align=PP_ALIGN.CENTER)
    _text(s, 0.55, cy + 2.0, 12.2, 0.8,
          [("WHY: passive video does not build hire-ready skills.", 12, True, SLATE),
           ("WHEN: one loop per study session (≈2 hrs). HOW: workbook steps → compare answer.", 11, False, GREY)],
          align=PP_ALIGN.CENTER)
    _footer(s, idx, FOOTER)
    return s


def s_phase_detail(prs, idx, acc, name, when, skills_text, id_from, id_to):
    s = _blank(prs)
    y = _header(s, name, kicker=f"{when} · Skills {id_from}–{id_to}")
    _rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, 0.55, y + 0.35, 12.2, 0.55, fill=acc)
    _text(s, 0.7, y + 0.45, 11.8, 0.4, [(skills_text, 12, True, CREAM)], align=PP_ALIGN.CENTER)
    cols = 2
    subset = [(j, sk) for j, sk in enumerate(SKILLS, 1) if id_from <= j <= id_to]
    cw, ch, gx, gy = 5.85, 1.05, 0.35, 0.18
    x0 = (SW - (cw * cols + gx)) / 2
    for k, (num, sk) in enumerate(subset):
        r, c = divmod(k, cols)
        x = x0 + c * (cw + gx)
        yy = y + 1.15 + r * (ch + gy)
        _rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, yy, cw, ch, fill=CARD_FILL, line=BORDER)
        _rect(s, MSO_SHAPE.RECTANGLE, x, yy, 0.12, ch, fill=acc)
        _text(s, x + 0.22, yy + 0.08, 0.4, 0.35, [(str(num), 11, True, acc)])
        _text(s, x + 0.65, yy + 0.06, cw - 0.75, 0.45, [(sk["name"], 10, True, INK)])
        _text(s, x + 0.65, yy + 0.52, cw - 0.75, 0.45, [(sk["exercise"][:55], 8.5, False, GREY)])
    _footer(s, idx, FOOTER)
    return s


def s_flow_diagram(prs, idx):
    s = _blank(prs)
    y = _header(s, "Portfolio build flow", kicker="What you ship by month 12")
    nodes = [
        (CLAY, "BRD\nrules script", 0.8, y + 0.6),
        (SKY, "SQL +\npandas KPIs", 3.2, y + 0.6),
        (OLIVE, "credit-pd\nmodel + SHAP", 5.6, y + 0.6),
        (FIG, "policy-rag\n+ agent", 8.0, y + 0.6),
        (CACTUS, "FastAPI\n+ Docker", 10.4, y + 0.6),
    ]
    nw, nh = 2.0, 1.1
    for acc, label, x, yy in nodes:
        _rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, yy, nw, nh, fill=acc)
        _text(s, x + 0.08, yy + 0.25, nw - 0.16, 0.6, [(label, 9, True, CREAM)], align=PP_ALIGN.CENTER)
    arrow_y = y + 1.15
    for i in range(len(nodes) - 1):
        x1 = nodes[i][2] + nw
        x2 = nodes[i + 1][2]
        conn = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, PInches(x1), PInches(arrow_y), PInches(x2), PInches(arrow_y))
        conn.line.color.rgb = GREY
    _rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, 3.5, y + 2.1, 6.5, 1.35, fill=LIGHT, line=OLIVE, line_w=1.5)
    _text(s, 3.7, y + 2.25, 6.1, 1.0,
          [("Month 11–12: eval harness · guardrails · CI · CV · demo video · apply OCB/NAB", 11, True, INK)],
          align=PP_ALIGN.CENTER)
    _footer(s, idx, FOOTER)
    return s


def s_workbook_pointer(prs, idx):
    s = _blank(prs)
    y = _header(s, "Study materials", kicker="Where to find detail")
    items = [
        (CLAY, "ai-skills-workbook.md", "Links · steps · answers · what/why/when/how"),
        (SKY, "AI-Skills-Learning-Slides.pptx", "19 skill slides · definition + exercise"),
        (OLIVE, "lab/exercises/", "week01 · week02 · solutions/"),
        (FIG, "Canvas: ai-skills-learning-path", "Interactive timeline · phase chart"),
        (CACTUS, "reading-path.md", "Book order · when to read each source"),
    ]
    for i, (acc, title, sub) in enumerate(items):
        yy = y + 0.4 + i * 0.95
        _rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, 1.2, yy, 11.0, 0.78, fill=CARD_FILL, line=BORDER)
        _rect(s, MSO_SHAPE.RECTANGLE, 1.2, yy, 0.14, 0.78, fill=acc)
        _text(s, 1.45, yy + 0.1, 10.5, 0.35, [(title, 12, True, INK)])
        _text(s, 1.45, yy + 0.42, 10.5, 0.3, [(sub, 10, False, GREY)])
    _footer(s, idx, FOOTER)
    return s


def generate():
    prs = Presentation()
    prs.slide_width = PInches(SW)
    prs.slide_height = PInches(SH)

    title_slide(
        _blank(prs),
        "AI Skills Visual Guide",
        "Learning Path at a Glance",
        [
            "Journey map · 4-step loop · 5 phases · portfolio flow",
            "Companion: ai-skills-workbook.md · Canvas in IDE",
        ],
    )

    slide_i = 2
    s_journey(prs, slide_i)
    slide_i += 1
    s_loop(prs, slide_i)
    slide_i += 1
    for acc, name, when, skills, id_from, id_to in PHASES:
        s_phase_detail(prs, slide_i, acc, name, when, skills, id_from, id_to)
        slide_i += 1
    s_flow_diagram(prs, slide_i)
    slide_i += 1
    s_workbook_pointer(prs, slide_i)

    summary_slide(
        _blank(prs),
        "Start today",
        "Skill 1 — Python",
        [
            (CLAY, "Read", "Python.org §1 + Kaggle lesson 1"),
            (SKY, "Do", "week01_brd_checklist.py"),
            (OLIVE, "Check", "Missing ACCEPTANCE CRITERIA in output"),
            (FIG, "Visual", "Open Canvas · ai-skills-learning-path"),
        ],
        "Workbook: curriculum/ai-skills-workbook.md",
    )

    out = OUTPUT / "AI-Skills-Visual-Slides.pptx"
    prs.save(out)
    print(f"Created {out} ({len(prs.slides._sldIdLst)} slides)")
    return out


if __name__ == "__main__":
    generate()
