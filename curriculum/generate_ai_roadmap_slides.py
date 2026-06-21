#!/usr/bin/env python3
"""Generate Zero-to-AI-Expert roadmap deck — Anthropic brand theme."""

import sys
from pathlib import Path
from repo_paths import EXPORTS_LEARNING

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches as PInches

ROOT = Path(__file__).resolve().parent
sys.path.insert(0, str(ROOT))
from anthropic_theme import (  # noqa: E402
    SW, SH, ACCENTS, CLAY, SKY, OLIVE, FIG, CACTUS,
    NAVY, BLUE, GREEN, ORANGE, PURPLE, TEAL, RED, WHITE, INK, GREY, LIGHT,
    BORDER, MID_GRAY, ON_DARK_MUTED, CREAM, CARD_FILL,
    blank as _blank, rect as _rect, text as _text, header as _header,
    footer as _footer, card as _card, title_slide, summary_slide,
)
from learning_data import META, TRACK_B_CHECKPOINTS  # noqa: E402

OUTPUT = EXPORTS_LEARNING
OUTPUT.mkdir(parents=True, exist_ok=True)

FOOTER_LABEL = "Zero to AI Expert  ·  Banking domain path"


def _footer_slide(slide, idx):
    _footer(slide, idx, FOOTER_LABEL)


def s_title(prs, i):
    s = _blank(prs)
    title_slide(
        s,
        "Learning Roadmap",
        "Zero to AI Expert",
        [
            "Banking domain · Start from zero Python · 12-month plan",
            f"{META['hours_per_week']} hrs/week ({META['hours_track_a']}h technical + {META['hours_track_b']}h Head of AI track)",
        ],
    )
    return s


def s_start_target(prs, i):
    s = _blank(prs)
    y = _header(s, "Where you are → where you're going", kicker="Starting point")
    cw, gap = 5.85, 0.45
    x0 = (SW - (cw * 2 + gap)) / 2
    _card(s, x0, y + 0.35, cw, 2.6, CLAY, "Today", "Banking domain · BRD/UAT · compliance mindset · Python = 0", "YOU")
    _card(s, x0 + cw + gap, y + 0.35, cw, 2.6, OLIVE, "Month 12 target", "2 banking AI projects · LangGraph copilot · ML model · HoAI artifacts (Track B)", "HIRE")
    _text(s, 0.55, y + 3.2, 12.2, 0.55,
          [("Expert in industry = job-ready portfolio + 2–3 yrs production — domain saves you ~1 year of \"what problem?\"", 13, True, NAVY)],
          align=PP_ALIGN.CENTER)
    _text(s, 0.55, y + 3.85, 12.2, 1.2,
          [("Learn every phase through banking exercises — never abstract-only tutorials", 12, False, GREY)],
          align=PP_ALIGN.CENTER)
    _footer_slide(s, i)
    return s


def s_gantt(prs, i):
    s = _blank(prs)
    y = _header(s, "12-month master timeline", kicker="At 10 hours per week")
    phases = [
        (SKY, "M1", "Python\n+ Git", 0.55, 0.9),
        (CACTUS, "M2", "SQL\n+ pandas", 1.55, 0.9),
        (OLIVE, "M3–4", "Classical\nML", 2.55, 1.8),
        (FIG, "M5–6", "NLP\n+ docs", 4.45, 1.8),
        (CLAY, "M7–8", "RAG\n+ LangGraph", 6.35, 1.8),
        (SKY, "M9–10", "FastAPI\n+ Docker", 8.25, 1.8),
        (OLIVE, "M11–12", "Portfolio\n+ apply", 10.15, 1.8),
    ]
    bar_y = y + 1.1
    _rect(s, MSO_SHAPE.RECTANGLE, 0.55, bar_y + 1.35, 11.5, 0.06, fill=MID_GRAY)
    for acc, when, label, x, w in phases:
        _rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, bar_y, w, 1.15, fill=acc)
        _text(s, x + 0.08, bar_y + 0.08, w - 0.16, 0.35, [(when, 11, True, CREAM)], align=PP_ALIGN.CENTER)
        _text(s, x + 0.08, bar_y + 0.42, w - 0.16, 0.65, [(label, 10, True, CREAM)], align=PP_ALIGN.CENTER)
    deliverables = [
        "BRD scorer", "KPI SQL", "PD model", "Doc classifier", "Policy copilot", "API + eval", "Job apps"
    ]
    dx = 0.55
    dw = 11.5 / len(deliverables)
    for j, d in enumerate(deliverables):
        _text(s, dx + j * dw, bar_y + 1.55, dw, 0.45, [(d, 9, False, INK)], align=PP_ALIGN.CENTER)
    _footer_slide(s, i)
    return s


def s_phase_cards(prs, i, title, kicker, items):
    s = _blank(prs)
    y = _header(s, title, kicker=kicker)
    n = len(items)
    cols = min(4, n)
    rows = (n + cols - 1) // cols
    cw, ch, gx, gy = 2.85, 1.55, 0.25, 0.28
    x0 = (SW - (cw * cols + gx * (cols - 1))) / 2
    for k, (acc, head, sub) in enumerate(items):
        r, c = divmod(k, cols)
        _card(s, x0 + c * (cw + gx), y + 0.35 + r * (ch + gy), cw, ch, acc, head, sub)
    _footer_slide(s, i)
    return s


def s_weekly(prs, i):
    s = _blank(prs)
    y = _header(s, "Weekly rhythm (10 hrs)", kicker="8h Track A · 2h Track B on milestone weeks")
    days = [
        (SKY, "Mon", "2h learn", "Course + notes"),
        (CACTUS, "Tue", "2h code", "Exercises"),
        (OLIVE, "Wed", "1h domain", "Banking exercise"),
        (FIG, "Thu", "2h project", "Deliverable"),
        (CLAY, "Fri", "1h ship", "Git + README"),
        (SKY, "Sat", "2h build", "Debug / Track B template"),
    ]
    cw, gap = 1.85, 0.22
    x0 = (SW - (cw * 6 + gap * 5)) / 2
    for j, (acc, day, hrs, act) in enumerate(days):
        x = x0 + j * (cw + gap)
        _rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, y + 0.5, cw, 2.2, fill=CARD_FILL, line=BORDER)
        _rect(s, MSO_SHAPE.RECTANGLE, x, y + 0.5, cw, 0.55, fill=acc)
        _text(s, x, y + 0.58, cw, 0.4, [(day, 13, True, CREAM)], align=PP_ALIGN.CENTER)
        _text(s, x + 0.12, y + 1.25, cw - 0.24, 0.45, [(hrs, 14, True, acc)], align=PP_ALIGN.CENTER)
        _text(s, x + 0.12, y + 1.75, cw - 0.24, 0.8, [(act, 10, False, GREY)], align=PP_ALIGN.CENTER)
    _text(s, 0.55, y + 3.0, 12.2, 0.5,
          [("Rule: never watch without coding the same day · Track B on W8/16/28/40/52 · Sun = rest", 12, True, NAVY)], align=PP_ALIGN.CENTER)
    _footer_slide(s, i)
    return s


def s_track_b(prs, i):
    s = _blank(prs)
    y = _header(s, "Track B — Head of AI Factory", kicker="2h/week · VPBank HoAI-aligned")
    for j, cp in enumerate(TRACK_B_CHECKPOINTS):
        col, row = j % 2, j // 2
        cw, ch, gx, gy = 5.85, 1.35, 0.45, 0.28
        x0 = (SW - (cw * 2 + gx)) / 2
        x = x0 + col * (cw + gx)
        yy = y + 0.35 + row * (ch + gy)
        acc = [CLAY, SKY, FIG, OLIVE, CACTUS][j]
        _rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, yy, cw, ch, fill=CARD_FILL, line=acc, line_w=2)
        _text(s, x + 0.2, yy + 0.15, 1.0, 0.35, [(cp["id"], 12, True, acc)])
        _text(s, x + 1.2, yy + 0.15, 4.5, 0.35, [(f"Week {cp['after_week']}", 10, False, GREY)])
        _text(s, x + 0.2, yy + 0.55, cw - 0.4, 0.7, [(cp["label"], 10, True, INK)])
    _text(s, 0.55, y + 3.35, 12.2, 0.45,
          [("Deck: exports/learning/Learning-Track-B-Slides.pptx · Guide: head-of-ai-track.md", 11, True, NAVY)],
          align=PP_ALIGN.CENTER)
    _footer_slide(s, i)
    return s


def s_domain_edge(prs, i):
    s = _blank(prs)
    y = _header(s, "Use banking domain as your edge", kicker="Project map")
    cases = [
        (FIG, "Credit PD model", "Classical ML · Month 4–5"),
        (SKY, "Policy RAG copilot", "LangGraph · Month 7–8"),
        (OLIVE, "Loan doc classifier", "NLP · Month 6"),
        (CLAY, "AML triage agent", "Optional · Month 11"),
        (CACTUS, "BRD quality scorer", "Python · Month 1"),
        (SKY, "Lending KPI SQL", "SQL · Month 2"),
    ]
    cols = 3
    cw, ch, gx, gy = 3.85, 1.45, 0.3, 0.3
    x0 = (SW - (cw * cols + gx * (cols - 1))) / 2
    for k, (acc, head, sub) in enumerate(cases):
        r, c = divmod(k, cols)
        _card(s, x0 + c * (cw + gx), y + 0.35 + r * (ch + gy), cw, ch, acc, head, sub)
    _footer_slide(s, i)
    return s


def s_milestones(prs, i):
    s = _blank(prs)
    y = _header(s, "Milestone checklist", kicker="Know when you're ready to apply")
    tiles = [
        (OLIVE, "✓", "Python + SQL", "Month 2"),
        (SKY, "✓", "ML + metrics", "Month 5"),
        (FIG, "✓", "RAG + LangGraph", "Month 8"),
        (CLAY, "✓", "API + Docker", "Month 10"),
        (CACTUS, "✓", "2 bank use cases", "Month 11"),
        (NAVY, "✓", "CV + GitHub", "Month 12"),
    ]
    cw, ch, gx = 3.85, 1.65, 0.35
    x0 = (SW - (cw * 3 + gx * 2)) / 2
    for k, (acc, arrow, head, sub) in enumerate(tiles):
        r, c = divmod(k, 3)
        x = x0 + c * (cw + gx)
        yy = y + 0.4 + r * (ch + 0.35)
        _rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, yy, cw, ch, fill=acc)
        _text(s, x + 0.25, yy + 0.15, cw - 0.5, 0.7, [(arrow, 28, True, CREAM)])
        _text(s, x + 0.25, yy + 0.85, cw - 0.5, 0.4, [(head, 13, True, CREAM)])
        _text(s, x + 0.25, yy + 1.25, cw - 0.5, 0.35, [(sub, 10, False, ON_DARK_MUTED)])
    _footer_slide(s, i)
    return s


def s_first30(prs, i):
    s = _blank(prs)
    y = _header(s, "First 30 days — start tomorrow", kicker="Action plan")
    weeks = [
        (SKY, "Week 1", "Install Python, Git, VS Code · Python tutorial §1–4 · GitHub repo"),
        (CACTUS, "Week 2", "Functions + files · BRD section checker (5 sections)"),
        (OLIVE, "Week 3", "SQLBolt 1–12 · SQLite + Python query"),
        (FIG, "Week 4", "pandas read_csv · one lending KPI chart"),
    ]
    cw, gap = 2.85, 0.3
    x0 = (SW - (cw * 4 + gap * 3)) / 2
    _rect(s, MSO_SHAPE.RECTANGLE, x0, y + 0.35, cw * 4 + gap * 3, 0.08, fill=MID_GRAY)
    for j, (acc, when, items) in enumerate(weeks):
        x = x0 + j * (cw + gap)
        _rect(s, MSO_SHAPE.OVAL, x + cw / 2 - 0.16, y + 0.27, 0.32, 0.32, fill=acc)
        _rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, y + 0.85, cw, 2.5, fill=CARD_FILL, line=BORDER)
        _rect(s, MSO_SHAPE.RECTANGLE, x, y + 0.85, cw, 0.65, fill=acc)
        _text(s, x, y + 0.92, cw, 0.5, [(when, 14, True, CREAM)], align=PP_ALIGN.CENTER)
        _text(s, x + 0.2, y + 1.65, cw - 0.4, 1.5, [(items, 10.5, False, INK)])
    _text(s, 0.55, y + 3.55, 12.2, 0.5,
          [("Full syllabus: curriculum/zero-to-ai-expert-syllabus.md", 12, True, NAVY)], align=PP_ALIGN.CENTER)
    _footer_slide(s, i)
    return s


def s_resources(prs, i):
    s = _blank(prs)
    y = _header(s, "What to learn — core stack", kicker="Syllabus summary")
    layers = [
        ("Month 1–2", "Python · Git · SQL · pandas", SKY),
        ("Month 3–5", "scikit-learn · XGBoost · SHAP", OLIVE),
        ("Month 5–6", "Embeddings · vector DB · NLP", CACTUS),
        ("Month 7–8", "LLM · RAG · LangGraph · eval", FIG),
        ("Month 9–10", "FastAPI · Docker · guardrails", CLAY),
        ("Always", "Banking metrics · governance · English", NAVY),
    ]
    lh, gap = 0.72, 0.14
    x, w = 1.2, SW - 2.4
    yy = y + 0.25
    for name, items, acc in layers:
        _rect(s, MSO_SHAPE.RECTANGLE, x, yy, 2.5, lh, fill=acc)
        _text(s, x + 0.15, yy, 2.2, lh, [(name, 11, True, CREAM)], anchor=MSO_ANCHOR.MIDDLE)
        _rect(s, MSO_SHAPE.RECTANGLE, x + 2.5, yy, w - 2.5, lh, fill=LIGHT, line=BORDER)
        _text(s, x + 2.65, yy, w - 2.7, lh, [(items, 11.5, False, INK)], anchor=MSO_ANCHOR.MIDDLE)
        yy += lh + gap
    _footer_slide(s, i)
    return s


def s_summary(prs, i):
    s = _blank(prs)
    summary_slide(
        s,
        "Summary",
        "Three rules for your path",
        [
            (SKY, "Learn through banking", "Every phase = BRD, credit, or ops artifact"),
            (OLIVE, "Ship, don't watch", "GitHub commit every week"),
            (CLAY, "Production beats PoC", "API + eval by month 10"),
        ],
        "Run: python3 curriculum/generate_ai_roadmap_slides.py  ·  Read: zero-to-ai-expert-syllabus.md",
    )
    return s


def generate():
    prs = Presentation()
    prs.slide_width = PInches(SW)
    prs.slide_height = PInches(SH)
    builders = [
        s_title,
        s_start_target,
        s_gantt,
        lambda p, i: s_phase_cards(p, i, "Phase 0–1 · Python & data (Months 1–3)",
                                    "Foundation",
                                    [(SKY, "Python syntax", "Variables · loops · functions"),
                                     (CACTUS, "Git + GitHub", "Commit every week"),
                                     (OLIVE, "SQL", "JOIN · GROUP BY · windows"),
                                     (FIG, "pandas", "CSV · groupby · plots"),
                                     (CLAY, "Deliverable", "BRD scorer + KPI queries"),
                                     (SKY, "Resources", "Python.org · SQLBolt · Kaggle")]),
        lambda p, i: s_phase_cards(p, i, "Phase 2 · Machine learning (Months 4–5)",
                                    "Classical ML",
                                    [(OLIVE, "scikit-learn", "Logistic · trees · RF"),
                                     (SKY, "XGBoost", "Credit PD model"),
                                     (CACTUS, "Metrics", "AUC · KS · calibration"),
                                     (FIG, "SHAP", "Explain declines"),
                                     (CLAY, "Deliverable", "credit-pd-model repo"),
                                     (FIG, "Domain tie", "BRD-style business metric")]),
        lambda p, i: s_phase_cards(p, i, "Phase 3–4 · GenAI & agents (Months 6–8)",
                                    "Hire-critical",
                                    [(FIG, "Embeddings", "Policy chunks"),
                                     (SKY, "RAG pipeline", "Retrieve + generate"),
                                     (CLAY, "LangGraph", "Tools · state · escalate"),
                                     (OLIVE, "Eval set", "20–50 golden Q&A"),
                                     (CACTUS, "Doc classifier", "Loan file types"),
                                     (FIG, "Deliverable", "policy-copilot-agent")]),
        lambda p, i: s_phase_cards(p, i, "Phase 5–7 · Production & job hunt (Months 9–12)",
                                    "Expert polish",
                                    [(CLAY, "FastAPI", "REST /ask endpoint"),
                                     (SKY, "Docker", "docker compose up"),
                                     (OLIVE, "Guardrails", "PII · audit trail"),
                                     (FIG, "Use case #2", "AML or NBO"),
                                     (CACTUS, "Portfolio", "GitHub + README"),
                                     (NAVY, "Apply", "NAB · OCB · VPBank · Anthropic")]),
        s_weekly,
        s_track_b,
        s_domain_edge,
        s_resources,
        s_milestones,
        s_first30,
        s_summary,
    ]
    for idx, b in enumerate(builders, 1):
        b(prs, idx)
    out = OUTPUT / "Zero-to-AI-Expert-Roadmap-Slides.pptx"
    prs.save(out)
    print(f"Created {out} ({len(prs.slides._sldIdLst)} slides)")
    return out


if __name__ == "__main__":
    generate()
