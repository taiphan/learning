#!/usr/bin/env python3
"""Track B + VPBank Head of AI Factory steering deck — from head-of-ai-track.md & vpbank one-pager."""

from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches as PInches

ROOT = Path(__file__).resolve().parent
sys.path.insert(0, str(ROOT))

from anthropic_theme import (  # noqa: E402
    SW,
    SH,
    CLAY,
    SKY,
    OLIVE,
    FIG,
    CACTUS,
    NAVY,
    INK,
    GREY,
    LIGHT,
    BORDER,
    CREAM,
    CARD_FILL,
    MID_GRAY,
    blank as _blank,
    rect as _rect,
    text as _text,
    header as _header,
    footer as _footer,
    card as _card,
    title_slide,
    summary_slide,
)
from learning_data import META, TRACK_B_CHECKPOINTS  # noqa: E402
from repo_paths import EXPORTS_LEARNING  # noqa: E402

OUTPUT = EXPORTS_LEARNING
OUTPUT.mkdir(parents=True, exist_ok=True)
FOOTER = "Track B · Head of AI Factory · curriculum/head-of-ai-track.md"


def _f(slide, idx):
    _footer(slide, idx, FOOTER)


def s_title(prs, idx):
    s = _blank(prs)
    title_slide(
        s,
        "Track B — Head of AI Factory",
        "Leadership alongside technical track",
        [
            f"{META['hours_track_a']}h/week ship code · {META['hours_track_b']}h/week strategy & governance",
            "Milestones: W8 · W16 · W28 · W40 · W52",
            "VPBank HoAI JD-aligned steering variant included",
        ],
    )
    return s


def s_dual_track(prs, idx):
    s = _blank(prs)
    y = _header(s, "Dual track · 10 hours/week", kicker="Track A proves you ship · Track B proves you lead")
    cw, gap = 5.85, 0.45
    x0 = (SW - (cw * 2 + gap)) / 2
    _card(s, x0, y + 0.4, cw, 2.8, SKY, "Track A (8h)", "Python · ML · RAG · agents · FastAPI · Docker", "SHIP")
    _card(s, x0 + cw + gap, y + 0.4, cw, 2.8, CLAY, "Track B (2h)", "Strategy · value case · governance · 90-day plan · steering deck", "LEAD")
    bullets = [
        "Same day as lab — leadership without code is slides-only",
        "Banking BA edge: BRD + compliance + KPI language in artifacts",
        "Year 1 target: AI Engineer — not Head of AI yet",
    ]
    body = "\n".join(f"• {b}" for b in bullets)
    _text(s, 1.0, y + 3.45, 11.5, 1.0, [(body, 11, False, INK)])
    _f(s, idx)
    return s


def s_milestones(prs, idx):
    s = _blank(prs)
    y = _header(s, "Track B milestones (H0–H4)", kicker="Mark complete in Learning app → Leadership tab")
    for i, cp in enumerate(TRACK_B_CHECKPOINTS):
        yy = y + 0.3 + i * 0.72
        _rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, 1.0, yy, 11.5, 0.58, fill=CARD_FILL, line=FIG if i % 2 == 0 else BORDER)
        _text(s, 1.15, yy + 0.12, 1.0, 0.35, [(cp["id"], 11, True, CLAY)])
        _text(s, 2.2, yy + 0.1, 1.8, 0.35, [(f"Week {cp['after_week']}", 10, False, GREY)])
        _text(s, 4.1, yy + 0.1, 8.2, 0.35, [(cp["label"], 10, True, INK)])
    _f(s, idx)
    return s


def s_vpbank_hook(prs, idx):
    s = _blank(prs)
    y = _header(s, "VPBank Head of AI Factory — steering hook", kicker="Scenario deck · job-skills-adaptation.md §F.4")
    _text(
        s,
        0.8,
        y + 0.25,
        11.8,
        0.9,
        [
            (
                "Value-led predictive + generative AI at production scale — "
                "governance-assured, subsidiary-ready.",
                14,
                True,
                NAVY,
            )
        ],
        align=PP_ALIGN.CENTER,
    )
    cards = [
        (OLIVE, "Smarter decisions", "Data-driven credit, ops, CX"),
        (SKY, "Lower risk", "PD, early warning, anomaly detection"),
        (CLAY, "Better experience", "Copilots with guardrails · human-in-loop"),
    ]
    cw, gap = 3.9, 0.35
    x0 = (SW - (cw * 3 + gap * 2)) / 2
    for i, (acc, head, sub) in enumerate(cards):
        _card(s, x0 + i * (cw + gap), y + 1.35, cw, 2.2, acc, head, sub)
    _text(
        s,
        0.8,
        y + 3.85,
        11.8,
        0.5,
        [("Apply Year 4–5 · Near-term: VPBank AI Engineer / Senior BA", 11, False, GREY)],
        align=PP_ALIGN.CENTER,
    )
    _f(s, idx)
    return s


def s_five_pillars(prs, idx):
    s = _blank(prs)
    y = _header(s, "Five pillars · Year 1 outcomes", kicker="Slide 2 — Strategy")
    pillars = [
        (CLAY, "Business value", "≥2 pilots with exec sponsor + baseline KPI"),
        (SKY, "Predictive AI", "1 prod model · AUC + champion/challenger"),
        (FIG, "Generative AI", "Policy copilot G1→G2 · grounded ≥90%"),
        (OLIVE, "Platform & reuse", "Shared RAG stack · 2+ squads reuse"),
        (CACTUS, "Governance", "100% pass intake gate · zero PII in logs"),
    ]
    for i, (acc, head, sub) in enumerate(pillars):
        yy = y + 0.25 + i * 0.68
        _rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, 1.0, yy, 11.5, 0.58, fill=CARD_FILL, line=acc, line_w=2)
        _rect(s, MSO_SHAPE.RECTANGLE, 1.0, yy, 0.12, 0.58, fill=acc)
        _text(s, 1.25, yy + 0.1, 3.2, 0.38, [(head, 11, True, acc)])
        _text(s, 4.5, yy + 0.1, 7.8, 0.38, [(sub, 10, False, INK)])
    _text(s, 1.0, y + 3.75, 11.5, 0.35,
          [("Principle: Value-led · platform-powered · governance-assured", 11, True, NAVY)])
    _f(s, idx)
    return s


def s_portfolio(prs, idx):
    s = _blank(prs)
    y = _header(s, "Portfolio wins + build vs buy", kicker="Slide 3 — Evidence")
    rows = [
        ("Credit PD model", "AUC · NPL link", "sklearn · SHAP · MLflow", "Build — proprietary features"),
        ("Policy RAG copilot", "Grounded ≥90%", "RAG · vector DB · eval", "Buy LLM + RAG layer"),
        ("Policy agent G2", "Tool success · audit", "LangGraph · guardrails", "Build workflow"),
        ("BRD intake gate", "Quality ≥80%", "apps/brd · checklist", "Build — factory intake"),
    ]
    lh = 0.62
    yy = y + 0.35
    headers = ["Project", "Metric", "Stack", "Build vs buy"]
    col_w = [3.0, 2.5, 3.2, 2.5]
    x = 0.8
    for j, h in enumerate(headers):
        _text(s, x + sum(col_w[:j]), yy, col_w[j], 0.3, [(h, 9, True, GREY)])
    yy += 0.35
    for row in rows:
        for j, cell in enumerate(row):
            _text(s, 0.8 + sum(col_w[:j]), yy, col_w[j], lh - 0.1, [(cell, 9.5, False, INK)])
        yy += lh
    _f(s, idx)
    return s


def s_kpis_governance(prs, idx):
    s = _blank(prs)
    y = _header(s, "KPIs · risk tiers · 90-day plan", kicker="Slide 4 — Govern & measure")
    kpis = [
        "Time-to-value (idea → prod): ≤90 days for G1 pilots",
        "Post-deploy ROI measured: 100% of prod use cases",
        "Governance pass rate: 100% at gate",
    ]
    body = "\n".join(f"• {k}" for k in kpis)
    _rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, 0.8, y + 0.3, 5.8, 2.0, fill=CARD_FILL, line=SKY)
    _text(s, 0.95, y + 0.38, 5.5, 0.28, [("KPIs (steering view)", 9, True, SKY)])
    _text(s, 0.95, y + 0.75, 5.5, 1.4, [(body, 10, False, INK)])

    tiers = "G1 — internal search · RAG eval\nG2 — RM copilot · human review\nG3 — credit decisions · SBV + override"
    _rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, 6.85, y + 0.3, 5.8, 2.0, fill=CARD_FILL, line=FIG)
    _text(s, 7.0, y + 0.38, 5.5, 0.28, [("Risk tiers G1 / G2 / G3", 9, True, FIG)])
    _text(s, 7.0, y + 0.75, 5.5, 1.4, [(tiers, 10, False, INK)])

    phases = "Q1 Foundation · Q2 Pilot · Q3 Scale MLOps · Q4 Steer + Year 2 roadmap"
    _rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, 0.8, y + 2.5, 11.85, 0.75, fill=LIGHT, line=BORDER)
    _text(s, 0.95, y + 2.65, 11.5, 0.45, [(phases, 10.5, True, INK)])
    _f(s, idx)
    return s


def s_ask_career(prs, idx):
    s = _blank(prs)
    y = _header(s, "The ask · career arc", kicker="Slide 5 + realistic path")
    asks = [
        ("Pilot sponsor", "Retail/credit BU for G1 copilot"),
        ("Headcount", "+2 FTE: ML engineer + LLMOps"),
        ("Platform budget", "Vector DB · inference · eval tooling"),
    ]
    for i, (head, sub) in enumerate(asks):
        yy = y + 0.35 + i * 0.72
        _rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, 0.8, yy, 5.5, 0.58, fill=CARD_FILL, line=CLAY)
        _text(s, 0.95, yy + 0.08, 5.2, 0.25, [(head, 10, True, CLAY)])
        _text(s, 0.95, yy + 0.32, 5.2, 0.22, [(sub, 9.5, False, INK)])

    arc = [
        ("Y1", "Dual track · engineer portfolio + HoAI artifacts"),
        ("Y2", "AI Engineer · 2–3 prod use cases"),
        ("Y3", "Tech Lead · squad 3–5"),
        ("Y4–5", "Head of AI Factory candidate"),
    ]
    for i, (when, label) in enumerate(arc):
        yy = y + 0.35 + i * 0.72
        _rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, 6.85, yy, 5.8, 0.58, fill=OLIVE if i == 3 else CARD_FILL, line=BORDER)
        _text(s, 7.0, yy + 0.12, 1.0, 0.35, [(when, 10, True, OLIVE if i == 3 else CLAY)])
        _text(s, 8.05, yy + 0.12, 4.4, 0.35, [(label, 9.5, False, CREAM if i == 3 else INK)])
    _f(s, idx)
    return s


def s_summary(prs, idx):
    s = _blank(prs)
    summary_slide(
        s,
        "Summary",
        "Track B deliverables",
        [
            (SKY, "Templates", "curriculum/templates/hoai/ — fill at each milestone week"),
            (OLIVE, "VPBank variant", "vpbank_steering_one_pager.md → this deck"),
            (CLAY, "Learning app", "Leadership tab on W8 · 16 · 28 · 40 · 52"),
        ],
        "Regenerate: python3 curriculum/generate_track_b_slides.py",
    )
    return s


def generate():
    prs = Presentation()
    prs.slide_width = PInches(SW)
    prs.slide_height = PInches(SH)
    builders = [
        s_title,
        s_dual_track,
        s_milestones,
        s_vpbank_hook,
        s_five_pillars,
        s_portfolio,
        s_kpis_governance,
        s_ask_career,
        s_summary,
    ]
    for idx, fn in enumerate(builders, 1):
        fn(prs, idx)
    out = OUTPUT / "Learning-Track-B-Slides.pptx"
    prs.save(out)
    print(f"Created {out} ({len(prs.slides._sldIdLst)} slides)")
    return out


if __name__ == "__main__":
    generate()
