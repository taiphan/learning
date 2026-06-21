#!/usr/bin/env python3
"""Master learning deck — 52 weeks × 3 slides, clickable index, resource links."""

from __future__ import annotations

import sys
from pathlib import Path
from repo_paths import EXPORTS_LEARNING

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches as PInches

ROOT = Path(__file__).resolve().parent
sys.path.insert(0, str(ROOT))

from anthropic_theme import (  # noqa: E402
    SW, SH, CLAY, SKY, OLIVE, FIG, CACTUS, CREAM, INK, GREY, LIGHT, BORDER, CARD_FILL, LINK,
    blank as _blank, rect as _rect, text as _text, header as _header, footer as _footer, title_slide,
)
from learning_data import CHECKPOINTS, META, PHASES, TRACK_B_CHECKPOINTS, TRACK_B_WEEKS, WEEKS  # noqa: E402
from learning_loader import phase_color  # noqa: E402
from week_slide_builders import (  # noqa: E402
    SLIDES_PER_WEEK,
    add_week_slides,
    first_slide_index_for_week,
)

OUTPUT = EXPORTS_LEARNING
OUTPUT.mkdir(parents=True, exist_ok=True)
FOOTER = f"Learning Master  ·  {SLIDES_PER_WEEK} slides/week  ·  Track B on W8/16/28/40/52"

FIRST_WEEK_SLIDE = 7  # 0-based: after title, nav, 4 index, phase overview
INDEX_SLIDE_START = 2

_INDEX_LINKS: list[tuple] = []
_BACK_LINKS: list[tuple] = []
_NAV_LINKS: list[tuple] = []


def _apply_deferred_links(prs: Presentation) -> None:
    for shape, target in _INDEX_LINKS + _BACK_LINKS + _NAV_LINKS:
        shape.click_action.target_slide = prs.slides[target]


def s_nav_help(prs, idx):
    s = _blank(prs)
    y = _header(s, "How to navigate this deck", kicker="Single source: learning_data.py")
    tips = [
        (CLAY, "Index slides", "Click any W## box → jump to that week's overview (3 slides each)"),
        (SKY, "Week slides", "Overview → Study → Lab — use Prev/Next within the week"),
        (FIG, "Track B (◆ weeks)", "W8 · W16 · W28 · W40 · W52 — Head of AI leadership (~2h)"),
        (OLIVE, "Workbook", "ai-skills-workbook.md — steps, answers, what/why/when"),
        (CACTUS, "Lab", "lab/ — run commands from the learning app Lab tab"),
    ]
    for i, (acc, head, sub) in enumerate(tips):
        yy = y + 0.35 + i * 0.72
        _rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, 1.0, yy, 11.5, 0.72, fill=CARD_FILL, line=acc, line_w=1.5)
        _text(s, 1.15, yy + 0.1, 11.2, 0.28, [(head, 11, True, acc)])
        _text(s, 1.15, yy + 0.38, 11.2, 0.28, [(sub, 10, False, GREY)])
    _footer(s, idx, FOOTER)
    return s


def s_index_quarter(prs, idx, q_label: str, week_range: range):
    s = _blank(prs)
    y = _header(s, f"Week index — {q_label}", kicker="Click a week to open overview slide · ◆ = Track B")
    cols, rows = 7, 2
    cw, ch, gx, gy = 1.55, 0.62, 0.18, 0.2
    x0 = (SW - (cw * cols + gx * (cols - 1))) / 2
    for i, wnum in enumerate(week_range):
        if wnum > 52:
            break
        r, c = divmod(i, cols)
        x = x0 + c * (cw + gx)
        yy = y + 0.35 + r * (ch + gy)
        w = next(x for x in WEEKS if x["week"] == wnum)
        acc = phase_color(w["phase"])
        sh = _rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, yy, cw, ch, fill=acc)
        marker = "◆ " if wnum in TRACK_B_WEEKS else ""
        _text(s, x + 0.06, yy + 0.06, cw - 0.12, 0.22, [(f"{marker}W{wnum:02d}", 10, True, CREAM)])
        title = w["title"][:18] + ("…" if len(w["title"]) > 18 else "")
        _text(s, x + 0.06, yy + 0.28, cw - 0.12, 0.3, [(title, 7.5, False, CREAM)])
        target = first_slide_index_for_week(wnum, FIRST_WEEK_SLIDE)
        _INDEX_LINKS.append((sh, target))
    _footer(s, idx, FOOTER)
    return s


def s_phase_overview(prs, idx):
    s = _blank(prs)
    y = _header(s, "Six phases · 52 weeks · 156 slides", kicker=META["data_source"])
    for i, p in enumerate(PHASES):
        yy = y + 0.3 + i * 0.72
        acc = phase_color(p["id"])
        _rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, 0.8, yy, 11.7, 0.62, fill=CARD_FILL, line=acc, line_w=2)
        _rect(s, MSO_SHAPE.RECTANGLE, 0.8, yy, 0.12, 0.62, fill=acc)
        _text(s, 1.05, yy + 0.08, 2.2, 0.45, [(f"Weeks {p['weeks']}", 10, True, acc)])
        _text(s, 3.3, yy + 0.08, 4.5, 0.45, [(p["name"], 11, True, INK)])
        _text(s, 8.0, yy + 0.08, 4.3, 0.45, [(p["theme"], 9.5, False, GREY)])
    _footer(s, idx, FOOTER)
    return s


def s_checkpoints(prs, idx):
    s = _blank(prs)
    y = _header(s, "Checkpoints — do not skip", kicker="Pass before accelerating")
    for i, cp in enumerate(CHECKPOINTS):
        yy = y + 0.35 + i * 0.72
        _rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, 1.0, yy, 11.5, 0.58, fill=CARD_FILL, line=BORDER)
        _text(s, 1.15, yy + 0.12, 1.2, 0.35, [(cp["id"], 11, True, CLAY)])
        _text(s, 2.4, yy + 0.1, 2.0, 0.35, [(f"After week {cp['after_week']}", 10, False, GREY)])
        _text(s, 4.5, yy + 0.1, 7.8, 0.35, [(cp["label"], 10, True, INK)])
    _footer(s, idx, FOOTER)
    return s


def s_track_b_checkpoints(prs, idx):
    s = _blank(prs)
    y = _header(s, "Track B — Head of AI milestones", kicker=f"{META['hours_track_b']}h/week · W8 · W16 · W28 · W40 · W52")
    for i, cp in enumerate(TRACK_B_CHECKPOINTS):
        yy = y + 0.35 + i * 0.72
        _rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, 1.0, yy, 11.5, 0.58, fill=CARD_FILL, line=CLAY if i == 0 else BORDER)
        _text(s, 1.15, yy + 0.12, 1.2, 0.35, [(cp["id"], 11, True, FIG)])
        _text(s, 2.4, yy + 0.1, 2.0, 0.35, [(f"After week {cp['after_week']}", 10, False, GREY)])
        _text(s, 4.5, yy + 0.1, 7.8, 0.35, [(cp["label"], 10, True, INK)])
    _text(s, 1.0, y + 4.0, 11.5, 0.45,
          [("Guide: curriculum/head-of-ai-track.md · VPBank deck: Learning-Track-B-Slides.pptx", 10, False, GREY)])
    _footer(s, idx, FOOTER)
    return s


def _index_slide_for_week(wnum: int) -> int:
    if wnum <= 13:
        return INDEX_SLIDE_START
    if wnum <= 26:
        return INDEX_SLIDE_START + 1
    if wnum <= 39:
        return INDEX_SLIDE_START + 2
    return INDEX_SLIDE_START + 3


def generate():
    global _INDEX_LINKS, _BACK_LINKS, _NAV_LINKS
    _INDEX_LINKS = []
    _BACK_LINKS = []
    _NAV_LINKS = []

    prs = Presentation()
    prs.slide_width = PInches(SW)
    prs.slide_height = PInches(SH)

    title_slide(
        _blank(prs),
        "52-Week Learning Master",
        "Banking domain → AI engineer · Head of AI (Track B)",
        [
            "Single data source: curriculum/learning_data.py",
            f"{SLIDES_PER_WEEK} slides/week · 156 week slides + index",
            f"{META['hours_per_week']} hrs/week ({META['hours_track_a']}h technical + {META['hours_track_b']}h leadership)",
        ],
    )

    s_nav_help(prs, 2)
    s_index_quarter(prs, 3, "Q1 · Weeks 1–13", range(1, 14))
    s_index_quarter(prs, 4, "Q2 · Weeks 14–26", range(14, 27))
    s_index_quarter(prs, 5, "Q3 · Weeks 27–39", range(27, 40))
    s_index_quarter(prs, 6, "Q4 · Weeks 40–52", range(40, 53))
    s_phase_overview(prs, 7)

    slide_num = 8
    for w in WEEKS:
        idx_slide = _index_slide_for_week(w["week"])
        add_week_slides(prs, w, idx_slide, _BACK_LINKS, _NAV_LINKS, slide_num)
        slide_num += SLIDES_PER_WEEK

    s_checkpoints(prs, slide_num)
    slide_num += 1
    s_track_b_checkpoints(prs, slide_num)
    _apply_deferred_links(prs)

    out = OUTPUT / "Learning-Master-Slides.pptx"
    prs.save(out)
    n = len(prs.slides._sldIdLst)
    print(f"Created {out} ({n} slides)")
    print(f"  Week 1 overview: slide {FIRST_WEEK_SLIDE + 1} (1-based)")
    print(f"  {SLIDES_PER_WEEK} slides per week × 52 = {52 * SLIDES_PER_WEEK} week slides")
    return out


if __name__ == "__main__":
    generate()
