"""Shared PPTX builders for per-week learning slides."""

from __future__ import annotations

import re
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches as PInches, Pt as PPt

from anthropic_theme import (
    SW,
    CLAY,
    SKY,
    OLIVE,
    FIG,
    CACTUS,
    CREAM,
    INK,
    GREY,
    LIGHT,
    BORDER,
    CARD_FILL,
    LINK,
    blank as _blank,
    rect as _rect,
    text as _text,
    header as _header,
    footer as _footer,
    title_slide,
)
from learning_loader import phase_color, skill_by_id

SLIDES_PER_WEEK = 3
FOOTER = "Learning  ·  learning_data.py  ·  Anthropic theme"

PHASE_TIPS: dict[str, str] = {
    "foundation": "Type every example yourself — no copy-paste. Commit when the script runs clean.",
    "ml": "Connect every metric to business language: NPL, false decline, approval rate.",
    "nlp": "Always log which source chunk grounded your answer.",
    "genai": "Refuse when context is missing — escalation is a feature, not failure.",
    "production": "If a friend cannot curl your API, it is not done.",
    "career": "Quantify impact in VND or bps — hiring managers need numbers.",
}


def study_bullets(study: str, max_items: int = 4) -> list[str]:
    for sep in (" · ", "; ", " ·", ";"):
        if sep in study:
            parts = [s.strip() for s in study.split(sep) if s.strip()]
            return parts[:max_items]
    return [study.strip()] if study.strip() else ["See workbook for reading list"]


def week_slug(week: dict) -> str:
    title = week["title"].lower()
    title = re.sub(r"[^a-z0-9]+", "-", title).strip("-")
    return f"week-{week['week']:02d}-{title[:40]}"


def _url_link(slide, x, y, w, h, label: str, url: str) -> None:
    if not url:
        return
    box = slide.shapes.add_textbox(PInches(x), PInches(y), PInches(w), PInches(h))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = label
    run.font.size = PPt(10)
    run.font.color.rgb = LINK
    run.hyperlink.address = url


def _code_box(slide, x, y, w, h, command: str) -> None:
    _rect(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, fill=LIGHT, line=BORDER)
    _text(slide, x + 0.15, y + 0.12, w - 0.3, h - 0.2, [(command, 10, False, INK)])


def _nav_back(slide, y_top: float, index_slide: int, back_links: list) -> None:
    back = _rect(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 10.2, y_top + 0.05, 2.5, 0.45, fill=LIGHT, line=BORDER)
    _text(slide, 10.3, y_top + 0.12, 2.3, 0.32, [("← Index", 9, True, INK)], align=PP_ALIGN.CENTER)
    back_links.append((back, index_slide))


def _nav_between(slide, y_top: float, wnum: int, part: int, prev_idx: int | None, next_idx: int | None, nav_links: list) -> None:
    label = f"Week {wnum} · {part}/{SLIDES_PER_WEEK}"
    _text(slide, 0.55, y_top + 0.12, 4.0, 0.32, [(label, 9, False, GREY)])
    if prev_idx is not None:
        prev = _rect(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 7.5, y_top + 0.05, 1.35, 0.45, fill=LIGHT, line=BORDER)
        _text(slide, 7.55, y_top + 0.12, 1.25, 0.32, [("← Prev", 8, True, INK)], align=PP_ALIGN.CENTER)
        nav_links.append((prev, prev_idx))
    if next_idx is not None:
        nxt = _rect(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 8.95, y_top + 0.05, 1.35, 0.45, fill=LIGHT, line=BORDER)
        _text(slide, 9.0, y_top + 0.12, 1.25, 0.32, [("Next →", 8, True, INK)], align=PP_ALIGN.CENTER)
        nav_links.append((nxt, next_idx))


def s_week_overview(
    prs: Presentation,
    idx: int,
    week: dict,
    index_slide: int,
    back_links: list,
    nav_links: list,
    slide_index_in_week: int,
    prev_slide: int | None,
    next_slide: int | None,
) -> Any:
    s = _blank(prs)
    wnum = week["week"]
    skill = skill_by_id(week["skill_id"])
    acc = phase_color(week["phase"])
    y = _header(s, f"Week {wnum} — {week['title']}", kicker=f"Month {week['month']} · {skill['name']}")

    _rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, 0.55, y + 0.2, 12.2, 0.55, fill=acc)
    _text(s, 0.7, y + 0.32, 11.8, 0.35, [(f"Deliverable: {week['deliverable']}", 12, True, CREAM)])

    _rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, 0.55, y + 0.95, 5.9, 2.2, fill=CARD_FILL, line=BORDER)
    _rect(s, MSO_SHAPE.RECTANGLE, 0.55, y + 0.95, 5.9, 0.38, fill=CLAY)
    _text(s, 0.7, y + 1.02, 5.6, 0.28, [("THIS WEEK YOU WILL", 9, True, CREAM)])
    bullets = study_bullets(week["study"])
    body = "\n".join(f"• {b}" for b in bullets)
    _text(s, 0.7, y + 1.45, 5.6, 1.55, [(body, 10, False, INK)])

    _rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, 6.65, y + 0.95, 6.1, 2.2, fill=CARD_FILL, line=BORDER)
    _rect(s, MSO_SHAPE.RECTANGLE, 6.65, y + 0.95, 6.1, 0.38, fill=SKY)
    _text(s, 6.8, y + 1.02, 5.8, 0.28, [("SKILL CONTEXT", 9, True, CREAM)])
    ctx = f"{skill['definition']}\n\nPhase tip: {PHASE_TIPS.get(week['phase'], '')}"
    _text(s, 6.8, y + 1.45, 5.8, 1.55, [(ctx, 9.5, False, INK)])

    if week.get("checkpoint"):
        _rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, 0.55, y + 3.35, 12.2, 0.55, fill=LIGHT, line=CACTUS, line_w=2)
        _text(s, 0.7, y + 3.48, 11.8, 0.35, [(f"Checkpoint {week['checkpoint']}", 11, True, INK)])

    tb = week.get("track_b")
    if tb:
        ty = y + (4.05 if week.get("checkpoint") else 3.35)
        _rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, 0.55, ty, 12.2, 0.72, fill=CARD_FILL, line=CLAY, line_w=2)
        _rect(s, MSO_SHAPE.RECTANGLE, 0.55, ty, 12.2, 0.32, fill=CLAY)
        cp = tb.get("hoai_checkpoint", "Track B")
        _text(s, 0.7, ty + 0.04, 11.8, 0.28, [(f"TRACK B · {cp} · ~{tb.get('hours', 2)}h leadership", 9, True, CREAM)])
        line = f"{tb.get('title', '')} — template: {tb.get('template', '')}"
        if len(line) > 95:
            line = line[:92] + "…"
        _text(s, 0.7, ty + 0.38, 11.8, 0.28, [(line, 9, False, INK)])

    _nav_between(s, y, wnum, slide_index_in_week, prev_slide, next_slide, nav_links)
    _nav_back(s, y, index_slide, back_links)
    _footer(s, idx, FOOTER)
    return s


def s_week_study(
    prs: Presentation,
    idx: int,
    week: dict,
    index_slide: int,
    back_links: list,
    nav_links: list,
    slide_index_in_week: int,
    prev_slide: int | None,
    next_slide: int | None,
) -> Any:
    s = _blank(prs)
    wnum = week["week"]
    skill = skill_by_id(week["skill_id"])
    acc = phase_color(week["phase"])
    y = _header(s, f"Week {wnum} — Study", kicker=f"Skill {skill['id']}: {skill['name']}")

    _rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, 0.55, y + 0.25, 12.2, 1.45, fill=CARD_FILL, line=BORDER)
    _rect(s, MSO_SHAPE.RECTANGLE, 0.55, y + 0.25, 12.2, 0.38, fill=acc)
    _text(s, 0.7, y + 0.32, 11.8, 0.28, [("READ & WATCH", 9, True, CREAM)])
    _text(s, 0.7, y + 0.72, 11.8, 0.9, [(week["study"], 11, False, INK)])

    _rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, 0.55, y + 1.85, 6.0, 1.55, fill=CARD_FILL, line=BORDER)
    _rect(s, MSO_SHAPE.RECTANGLE, 0.55, y + 1.85, 6.0, 0.35, fill=OLIVE)
    _text(s, 0.7, y + 1.92, 5.7, 0.28, [("SOURCES (skill)", 9, True, CREAM)])
    _text(s, 0.7, y + 2.32, 5.7, 0.95, [(skill["sources"], 9.5, False, INK)])

    _rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, 6.75, y + 1.85, 6.0, 1.55, fill=CARD_FILL, line=BORDER)
    _rect(s, MSO_SHAPE.RECTANGLE, 6.75, y + 1.85, 6.0, 0.35, fill=FIG)
    _text(s, 6.9, y + 1.92, 5.7, 0.28, [("PRACTICE HABIT", 9, True, CREAM)])
    _text(s, 6.9, y + 2.32, 5.7, 0.95, [(skill["practice"], 9.5, False, INK)])

    link_y = y + 3.55
    _text(s, 0.55, link_y, 12.2, 0.28, [("Click resources:", 10, True, INK)])
    urls = week.get("resource_urls") or skill.get("resource_urls", [])
    for j, (label, url) in enumerate(urls[:4]):
        col, row = j % 2, j // 2
        _url_link(s, 0.55 + col * 6.2, link_y + 0.32 + row * 0.38, 6.0, 0.35, label, url)

    _nav_between(s, y, wnum, slide_index_in_week, prev_slide, next_slide, nav_links)
    _nav_back(s, y, index_slide, back_links)
    _footer(s, idx, FOOTER)
    return s


def s_week_lab(
    prs: Presentation,
    idx: int,
    week: dict,
    index_slide: int,
    back_links: list,
    nav_links: list,
    slide_index_in_week: int,
    prev_slide: int | None,
    next_slide: int | None,
) -> Any:
    s = _blank(prs)
    wnum = week["week"]
    skill = skill_by_id(week["skill_id"])
    acc = phase_color(week["phase"])
    y = _header(s, f"Week {wnum} — Lab", kicker="Exercise · Run · Prove")
    tb = week.get("track_b")

    _rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, 0.55, y + 0.25, 12.2, 1.15, fill=CARD_FILL, line=BORDER)
    _rect(s, MSO_SHAPE.RECTANGLE, 0.55, y + 0.25, 12.2, 0.38, fill=SKY)
    _text(s, 0.7, y + 0.32, 11.8, 0.28, [("EXERCISE FILE", 9, True, CREAM)])
    _text(s, 0.7, y + 0.72, 11.8, 0.6, [(week["exercise"], 10, True, INK)])

    _text(s, 0.55, y + 1.55, 12.2, 0.28, [("Run command:", 10, True, INK)])
    _code_box(s, 0.55, y + 1.85, 12.2, 0.75, week["run"])

    _rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, 0.55, y + 2.75, 12.2, 1.05 if not tb else 1.35, fill=CARD_FILL, line=acc, line_w=2)
    _text(s, 0.7, y + 2.88, 11.8, 0.28, [("DONE WHEN", 9, True, acc)])
    steps = [
        f"1. Open `{week['exercise']}` in Cursor",
        f"2. Run: `{week['run']}`",
        f"3. Output matches: {week['deliverable']}",
        f"4. Skill checkpoint: {skill['checkpoint']}",
    ]
    if tb:
        steps.extend([
            f"5. Track B ({tb.get('hoai_checkpoint', 'HoAI')}): {tb.get('action', '')}",
            f"6. Leadership done when: {tb.get('deliverable', '')}",
        ])
    _text(s, 0.7, y + 3.22, 11.8, 0.5, [("\n".join(steps), 9.5 if not tb else 8.5, False, INK)])

    _nav_between(s, y, wnum, slide_index_in_week, prev_slide, next_slide, nav_links)
    _nav_back(s, y, index_slide, back_links)
    _footer(s, idx, FOOTER)
    return s


def add_week_slides(
    prs: Presentation,
    week: dict,
    index_slide: int,
    back_links: list,
    nav_links: list,
    start_idx: int,
) -> int:
    """Add 3 slides for one week. Returns next free slide index (1-based count helper)."""
    wnum = week["week"]
    base = len(prs.slides._sldIdLst)
    over_idx = base
    study_idx = base + 1
    lab_idx = base + 2

    s_week_overview(prs, start_idx, week, index_slide, back_links, nav_links, 1, None, study_idx)
    s_week_study(prs, start_idx + 1, week, index_slide, back_links, nav_links, 2, over_idx, lab_idx)
    s_week_lab(prs, start_idx + 2, week, index_slide, back_links, nav_links, 3, study_idx, None)
    return start_idx + SLIDES_PER_WEEK


def first_slide_index_for_week(wnum: int, first_week_slide: int) -> int:
    return first_week_slide + (wnum - 1) * SLIDES_PER_WEEK


def build_single_week_deck(week: dict, out_path) -> None:
    """Mini deck: title + 3 week slides."""
    prs = Presentation()
    from pptx.util import Inches as PInches

    prs.slide_width = PInches(SW)
    prs.slide_height = PInches(7.5)

    skill = skill_by_id(week["skill_id"])
    title_slide(
        _blank(prs),
        f"Week {week['week']} · Month {week['month']}",
        week["title"],
        [
            skill["name"],
            week["deliverable"],
            f"Exercise: {week['exercise']}",
        ],
    )

    back_links: list = []
    nav_links: list = []
    add_week_slides(prs, week, 0, back_links, nav_links, 2)

    for shape, target in back_links + nav_links:
        shape.click_action.target_slide = prs.slides[target]

    prs.save(out_path)
