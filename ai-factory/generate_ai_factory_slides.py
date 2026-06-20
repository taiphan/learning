#!/usr/bin/env python3
"""Generate the Head of AI Factory deck (PPTX) — Anthropic brand theme, visual-first."""

import math
import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches as PInches, Pt as PPt

ROOT = Path(__file__).resolve().parent
sys.path.insert(0, str(ROOT))
from anthropic_theme import (  # noqa: E402
    SW, SH, ACCENTS, SLATE, CREAM, CLAY, SKY, OLIVE, FIG, CACTUS,
    NAVY, BLUE, GREEN, ORANGE, PURPLE, TEAL, RED, INK, GREY, LIGHT,
    BORDER, LINK, ON_DARK_MUTED, SLATE_SOFT, CARD_FILL, MID_GRAY,
    blank as _blank, rect as _rect, text as _text, set_shape_text as _set_shape_text,
    header as _header, footer as _footer, card as _card,
    title_slide, summary_slide, slide_bg,
)

OUTPUT = ROOT / "exports"
OUTPUT.mkdir(exist_ok=True)

FOOTER_LABEL = "Head of AI Factory  ·  Board briefing"

def _footer_slide(slide, idx):
    _footer(slide, idx, FOOTER_LABEL)


# ---------- Slide builders ----------

def s_title(prs, idx):
    s = _blank(prs)
    title_slide(
        s,
        "Enterprise Data Science & AI",
        "Head of AI Factory",
        ["Turning AI opportunity into trusted, scaled business value"],
        "Predictive AI  ·  Generative AI  ·  Platform  ·  Governance  ·  People",
    )
    return s


def s_purpose(prs, idx):
    s = _blank(prs)
    y = _header(s, "Why this role exists", kicker="Role purpose")
    _text(s, 0.55, y, 12.2, 0.5, [("Lead AI & analytics that move three business needles", 15, False, GREY)])
    cards = [
        (BLUE, "Smarter decisions", "Data-driven choices across the enterprise", "01"),
        (GREEN, "Lower risk", "Better risk management and early warning", "02"),
        (ORANGE, "Better experience", "Personalized, faster customer journeys", "03"),
    ]
    cw, gap = 3.9, 0.35
    x0 = (SW - (cw * 3 + gap * 2)) / 2
    for i, (acc, head, sub, big) in enumerate(cards):
        _card(s, x0 + i * (cw + gap), y + 0.7, cw, 3.0, acc, head, sub, big)
    _text(s, 0.55, y + 4.05, 12.2, 0.6,
          [("Trusted advisor to executives — drives a culture of innovation and data-driven decisions", 13, True, NAVY)],
          align=PP_ALIGN.CENTER)
    _footer_slide(s, idx)
    return s


def s_hub_spoke(prs, idx):
    s = _blank(prs)
    y = _header(s, "What the AI Factory owns", kicker="Mandate & scope")
    cx, cy = SW / 2, 4.35
    hub = _rect(s, MSO_SHAPE.OVAL, cx - 1.15, cy - 1.15, 2.3, 2.3, fill=SLATE, line=CREAM, line_w=2)
    _set_shape_text(hub, [("AI", 26, True, CREAM), ("FACTORY", 14, True, CLAY)])
    spokes = [
        ("Predictive AI", "forecasting · scoring", BLUE),
        ("Generative AI", "LLMs · RAG · agents", PURPLE),
        ("Platform", "reusable assets · infra", TEAL),
        ("Governance", "security · ethics", RED),
        ("People", "talent · culture", GREEN),
    ]
    r = 2.55
    sw_, sh_ = 2.55, 0.95
    n = len(spokes)
    for i, (head, sub, acc) in enumerate(spokes):
        ang = -math.pi / 2 + i * (2 * math.pi / n)
        sx = cx + r * math.cos(ang) * 1.55
        sy = cy + r * math.sin(ang)
        conn = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                      PInches(cx), PInches(cy), PInches(sx), PInches(sy))
        conn.line.color.rgb = LINK
        conn.line.width = PPt(2)
        box = _rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, sx - sw_ / 2, sy - sh_ / 2, sw_, sh_, fill=acc, line=CREAM, line_w=1.5)
        _set_shape_text(box, [(head, 13, True, CREAM), (sub, 9.5, False, ON_DARK_MUTED)])
    _footer_slide(s, idx)
    return s


def s_loop(prs, idx):
    s = _blank(prs)
    y = _header(s, "How the AI Factory runs", kicker="Operating model")
    _text(s, 0.55, y, 12.2, 0.45, [("A repeatable value loop — reuse platform, governance & MLOps every cycle", 14, False, GREY)])
    nodes = [
        ("Intake\nuse-case", BLUE),
        ("Prioritize\nby value", SKY),
        ("Build &\nvalidate", GREEN),
        ("Governance\ngate", RED),
        ("Deploy\n(MLOps)", ORANGE),
        ("Monitor\nvalue", PURPLE),
    ]
    cx, cy = SW / 2, 4.5
    rx, ry = 4.5, 1.85
    n = len(nodes)
    pts = []
    for i in range(n):
        ang = -math.pi / 2 + i * (2 * math.pi / n)
        pts.append((cx + rx * math.cos(ang), cy + ry * math.sin(ang)))
    # arrows along the ring
    for i in range(n):
        x1, y1 = pts[i]
        x2, y2 = pts[(i + 1) % n]
        conn = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, PInches(x1), PInches(y1), PInches(x2), PInches(y2))
        conn.line.color.rgb = MID_GRAY
        conn.line.width = PPt(2.25)
    for i, (label, acc) in enumerate(nodes):
        x1, y1 = pts[i]
        d = 1.5
        node = _rect(s, MSO_SHAPE.OVAL, x1 - d / 2, y1 - d / 2, d, d, fill=acc, line=CREAM, line_w=2)
        _set_shape_text(node, [(label, 11, True, CREAM)])
    _footer_slide(s, idx)
    return s


def s_responsibilities(prs, idx):
    s = _blank(prs)
    y = _header(s, "Areas of responsibility", kicker="Accountability")
    items = [
        ("Strategy", "Enterprise DS & AI strategy"),
        ("Advisory", "Advise C-level on growth & risk"),
        ("Delivery", "Predictive + Generative AI"),
        ("Architecture", "MLOps/LLMOps & lifecycle"),
        ("Platform", "Reusable assets & infra"),
        ("Partnership", "IT & Data Architecture"),
        ("Governance", "Security, privacy, ethics"),
        ("Value", "Use cases, ROI, outcomes"),
        ("People", "Talent & capability"),
        ("Culture", "Develop future AI leaders"),
    ]
    cols, rows = 5, 2
    cw, ch, gx, gy = 2.32, 1.75, 0.16, 0.25
    x0 = (SW - (cw * cols + gx * (cols - 1))) / 2
    for i, (head, sub) in enumerate(items):
        r, c = divmod(i, cols)
        x = x0 + c * (cw + gx)
        yy = y + 0.35 + r * (ch + gy)
        acc = ACCENTS[i % len(ACCENTS)]
        _rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, yy, cw, ch, fill=CARD_FILL, line=BORDER, line_w=1)
        num = _rect(s, MSO_SHAPE.OVAL, x + 0.18, yy + 0.18, 0.5, 0.5, fill=acc)
        _set_shape_text(num, [(str(i + 1), 13, True, CREAM)])
        _text(s, x + 0.8, yy + 0.2, cw - 0.95, 0.45, [(head, 12.5, True, INK)])
        _text(s, x + 0.22, yy + 0.85, cw - 0.4, 0.8, [(sub, 10, False, GREY)])
    _footer_slide(s, idx)
    return s


def s_org(prs, idx):
    s = _blank(prs)
    y = _header(s, "AI Factory organization", kicker="Team & resources")
    root_w, root_h = 4.2, 0.75
    rx = (SW - root_w) / 2
    rb = _rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, rx, y + 0.15, root_w, root_h, fill=SLATE, line=CREAM, line_w=1.5)
    _set_shape_text(rb, [("Head of AI Factory", 14, True, CREAM)])
    groups = [
        ("Data Science / ML", ["Data Scientists", "ML Engineers"], BLUE),
        ("Generative AI", ["LLM / NLP Eng", "Prompt / RAG Eng"], PURPLE),
        ("Platform / MLOps", ["MLOps / DevOps", "Data Eng liaison"], TEAL),
        ("Governance & Risk", ["Model Validation", "Responsible AI"], RED),
        ("Business / Value", ["Use-case owners", "ROI & adoption"], GREEN),
    ]
    n = len(groups)
    gap = 0.22
    cw = (SW - 0.6 - gap * (n - 1)) / n
    x0 = 0.3
    mgr_top = y + 1.35
    mgr_h = 0.75
    root_cx = rx + root_w / 2
    for i, (mgr, reports, acc) in enumerate(groups):
        cx = x0 + i * (cw + gap)
        mcx = cx + cw / 2
        conn = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, PInches(root_cx), PInches(y + 0.15 + root_h),
                                      PInches(mcx), PInches(mgr_top))
        conn.line.color.rgb = LINK
        conn.line.width = PPt(1.25)
        mb = _rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, cx, mgr_top, cw, mgr_h, fill=acc, line=CREAM, line_w=1.5)
        _set_shape_text(mb, [(mgr, 11, True, CREAM)])
        ry = mgr_top + mgr_h + 0.3
        for rep in reports:
            rb2 = _rect(s, MSO_SHAPE.RECTANGLE, cx + 0.12, ry, cw - 0.24, 0.6, fill=LIGHT, line=acc, line_w=1)
            _set_shape_text(rb2, [(rep, 9.5, False, INK)])
            ry += 0.78
    _footer_slide(s, idx)
    return s


def s_portfolio(prs, idx):
    s = _blank(prs)
    y = _header(s, "Where AI creates value", kicker="Use-case portfolio (illustrative)")
    cards = [
        (RED, "Credit & Risk", "Scoring · PD/LGD · early warning · policy copilots"),
        (GREEN, "Growth & CX", "Propensity · churn · NBO · service agents"),
        (BLUE, "Operations", "Forecasting · optimization · knowledge RAG"),
        (PURPLE, "Fraud", "Anomaly detection · investigation assistants"),
        (ORANGE, "Finance", "Cashflow forecast · report-gen copilots"),
    ]
    cw, gap = 3.95, 0.3
    # row 1: 3 cards, row 2: 2 cards centered
    x0 = (SW - (cw * 3 + gap * 2)) / 2
    for i in range(3):
        acc, h, sub = cards[i]
        _card(s, x0 + i * (cw + gap), y + 0.35, cw, 1.85, acc, h, sub)
    x1 = (SW - (cw * 2 + gap)) / 2
    for j in range(2):
        acc, h, sub = cards[3 + j]
        _card(s, x1 + j * (cw + gap), y + 2.45, cw, 1.85, acc, h, sub)
    _text(s, 0.55, y + 4.45, 12.2, 0.4,
          [("Selection: business value · data readiness · feasibility · regulatory · reuse", 11.5, True, NAVY)],
          align=PP_ALIGN.CENTER)
    _footer_slide(s, idx)
    return s


def s_stack(prs, idx):
    s = _blank(prs)
    y = _header(s, "Technology stack", kicker="Platform foundations")
    layers = [
        ("Languages", "Python · R · SQL", BLUE),
        ("ML / Deep Learning", "scikit-learn · XGBoost · PyTorch · TensorFlow", SKY),
        ("Generative AI", "LLMs · RAG · vector DBs · agent frameworks", PURPLE),
        ("Distributed compute", "Spark · Hadoop", TEAL),
        ("Cloud (multi-cloud)", "AWS · GCP · Azure", GREEN),
        ("MLOps / LLMOps", "pipelines · registry · eval · guardrails · monitoring", ORANGE),
        ("Data foundation", "lakehouse · feature store · lineage", NAVY),
    ]
    lh = 0.66
    gap = 0.12
    x = 1.4
    w = SW - 2 * x
    yy = y + 0.2
    for name, items, acc in layers:
        _rect(s, MSO_SHAPE.RECTANGLE, x, yy, 3.4, lh, fill=acc)
        _text(s, x + 0.2, yy, 3.0, lh, [(name, 12, True, CREAM)], anchor=MSO_ANCHOR.MIDDLE)
        _rect(s, MSO_SHAPE.RECTANGLE, x + 3.4, yy, w - 3.4, lh, fill=LIGHT, line=BORDER, line_w=1)
        _text(s, x + 3.6, yy, w - 3.6, lh, [(items, 11.5, False, INK)], anchor=MSO_ANCHOR.MIDDLE)
        yy += lh + gap
    _footer_slide(s, idx)
    return s


def _chevron_flow(prs, idx, title, kicker, steps, note):
    s = _blank(prs)
    y = _header(s, title, kicker=kicker)
    n = len(steps)
    gap = 0.12
    cw = (SW - 1.0 - gap * (n - 1)) / n
    ch = 1.7
    x0 = 0.5
    yy = y + 0.85
    for i, label in enumerate(steps):
        acc = ACCENTS[i % len(ACCENTS)]
        sp = _rect(s, MSO_SHAPE.CHEVRON, x0 + i * (cw + gap), yy, cw, ch, fill=acc, line=CREAM, line_w=1.5)
        _set_shape_text(sp, [(f"{i+1}", 16, True, CREAM), (label, 11.5, True, CREAM)])
    _text(s, 0.55, yy + ch + 0.5, 12.2, 0.8, [(note, 13, False, GREY)], align=PP_ALIGN.CENTER)
    _footer_slide(s, idx)
    return s


def s_gates(prs, idx):
    s = _blank(prs)
    y = _header(s, "AI governance — three gates", kicker="Trust by design")
    gates = [
        (BLUE, "Gate 1 — Intake", "Before build", ["Business value", "Data availability", "Legal basis", "Risk tier"]),
        (ORANGE, "Gate 2 — Pre-deploy", "Before production", ["Validation", "Risk & ethics", "Security & privacy"]),
        (GREEN, "Gate 3 — In-life", "Continuous", ["Drift monitoring", "Performance", "Audit trail"]),
    ]
    cw, gap = 3.95, 0.35
    x0 = (SW - (cw * 3 + gap * 2)) / 2
    for i, (acc, head, when, checks) in enumerate(gates):
        x = x0 + i * (cw + gap)
        _rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, y + 0.3, cw, 3.5, fill=CARD_FILL, line=BORDER, line_w=1)
        _rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, y + 0.3, cw, 0.7, fill=acc)
        _set_shape_text(_rect(s, MSO_SHAPE.RECTANGLE, x, y + 0.3, cw, 0.7, fill=acc), [(head, 13, True, CREAM)])
        _text(s, x + 0.25, y + 1.1, cw - 0.5, 0.35, [(when.upper(), 10, True, acc)])
        for j, c in enumerate(checks):
            cyt = y + 1.55 + j * 0.5
            chk = _rect(s, MSO_SHAPE.OVAL, x + 0.28, cyt + 0.03, 0.26, 0.26, fill=acc)
            _set_shape_text(chk, [("✓", 10, True, CREAM)])
            _text(s, x + 0.68, cyt, cw - 0.9, 0.45, [(c, 11.5, False, INK)])
        # arrow between gates
        if i < 2:
            ar = _rect(s, MSO_SHAPE.RIGHT_ARROW, x + cw + 0.02, y + 1.7, gap - 0.04, 0.5, fill=LINK)
    _text(s, 0.55, y + 4.0, 12.2, 0.4,
          [("Risk tiers: Low · Medium · High · Critical → deeper control at higher tiers", 11.5, True, NAVY)],
          align=PP_ALIGN.CENTER)
    _footer_slide(s, idx)
    return s


def s_responsible(prs, idx):
    s = _blank(prs)
    y = _header(s, "Responsible AI", kicker="Principles in practice")
    tiles = [
        (BLUE, "Fairness", "Bias testing + mitigation"),
        (TEAL, "Transparency", "Model cards · explainability"),
        (NAVY, "Accountability", "Named business + model owner"),
        (GREEN, "Privacy", "Minimization · consent · DLP"),
        (RED, "Security", "Access · adversarial robustness"),
        (PURPLE, "Human oversight", "Human-in-the-loop high/critical"),
    ]
    cols = 3
    cw, ch, gx, gy = 3.85, 1.55, 0.3, 0.3
    x0 = (SW - (cw * cols + gx * (cols - 1))) / 2
    for i, (acc, head, sub) in enumerate(tiles):
        r, c = divmod(i, cols)
        x = x0 + c * (cw + gx)
        yy = y + 0.35 + r * (ch + gy)
        _rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, yy, cw, ch, fill=CARD_FILL, line=BORDER, line_w=1)
        _rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, yy, cw, 0.12, fill=acc)
        _text(s, x + 0.3, yy + 0.25, cw - 0.5, 0.5, [(head, 14, True, acc)])
        _text(s, x + 0.3, yy + 0.8, cw - 0.5, 0.6, [(sub, 11, False, GREY)])
    _footer_slide(s, idx)
    return s


def s_kpi(prs, idx):
    s = _blank(prs)
    y = _header(s, "How we measure success", kicker="KPI dashboard")
    tiles = [
        (GREEN, "ROI", "Net AI portfolio value", "↑"),
        (BLUE, "Delivery", "Use cases live / quarter", "↑"),
        (ORANGE, "Time-to-value", "Idea → production", "↓"),
        (TEAL, "Reuse", "% on shared platform", "↑"),
        (PURPLE, "GenAI safety", "Grounded-response rate", "↑"),
        (NAVY, "Governance", "Gate first-pass rate", "↑"),
        (SKY, "Quality", "vs baseline; drift stable", "→"),
        (RED, "Talent", "Key-talent attrition", "↓"),
    ]
    cols = 4
    cw, ch, gx, gy = 2.85, 1.85, 0.25, 0.3
    x0 = (SW - (cw * cols + gx * (cols - 1))) / 2
    for i, (acc, head, sub, arrow) in enumerate(tiles):
        r, c = divmod(i, cols)
        x = x0 + c * (cw + gx)
        yy = y + 0.3 + r * (ch + gy)
        _rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, yy, cw, ch, fill=acc)
        _text(s, x + 0.25, yy + 0.18, cw - 0.5, 0.85, [(arrow, 30, True, CREAM)])
        _text(s, x + 0.25, yy + 0.95, cw - 0.5, 0.4, [(head, 14, True, CREAM)])
        _text(s, x + 0.25, yy + 1.35, cw - 0.5, 0.45, [(sub, 10, False, ON_DARK_MUTED)])
    _footer_slide(s, idx)
    return s


def s_requirements(prs, idx):
    s = _blank(prs)
    y = _header(s, "Who we are looking for", kicker="Requirements / profile")
    cards = [
        (BLUE, "Education", "Master's / Ph.D. — Data Science, CS, Statistics, Math"),
        (GREEN, "Experience", "10+ yrs DS/analytics · 5+ yrs leadership · large-scale production"),
        (PURPLE, "Technical", "ML · DL · NLP · big data · Python/R/SQL · Spark/Hadoop · AWS/GCP/Azure"),
        (ORANGE, "Leadership", "Stakeholder management · influence at executive level"),
    ]
    cw, ch, gx, gy = 5.9, 1.75, 0.4, 0.35
    x0 = (SW - (cw * 2 + gx)) / 2
    for i, (acc, head, sub) in enumerate(cards):
        r, c = divmod(i, 2)
        _card(s, x0 + c * (cw + gx), y + 0.4 + r * (ch + gy), cw, ch, acc, head, sub)
    _footer_slide(s, idx)
    return s


def s_90days(prs, idx):
    s = _blank(prs)
    y = _header(s, "First 90 days", kicker="Execution plan")
    phases = [
        (BLUE, "0–30 days", "Listen & assess", ["Maturity & data audit", "Meet executives & BUs", "Quick-win shortlist"]),
        (ORANGE, "30–60 days", "Frame & fund", ["Ratify AI strategy", "Governance baseline", "Prioritized portfolio"]),
        (GREEN, "60–90 days", "Deliver & scale", ["Ship 1–2 quick wins", "Stand up MLOps", "Operating model live"]),
    ]
    cw, gap = 3.95, 0.35
    x0 = (SW - (cw * 3 + gap * 2)) / 2
    # timeline bar
    _rect(s, MSO_SHAPE.RECTANGLE, x0, y + 0.35, cw * 3 + gap * 2, 0.08, fill=MID_GRAY)
    for i, (acc, when, head, items) in enumerate(phases):
        x = x0 + i * (cw + gap)
        dot = _rect(s, MSO_SHAPE.OVAL, x + cw / 2 - 0.16, y + 0.27, 0.32, 0.32, fill=acc, line=CREAM, line_w=2)
        _rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, y + 0.85, cw, 3.0, fill=CARD_FILL, line=BORDER, line_w=1)
        _rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, y + 0.85, cw, 0.75, fill=acc)
        _set_shape_text(_rect(s, MSO_SHAPE.RECTANGLE, x, y + 0.85, cw, 0.75, fill=acc),
                        [(when, 15, True, CREAM), (head, 11, True, ON_DARK_MUTED)])
        for j, it in enumerate(items):
            cyt = y + 1.8 + j * 0.6
            bullet = _rect(s, MSO_SHAPE.OVAL, x + 0.3, cyt + 0.06, 0.18, 0.18, fill=acc)
            _text(s, x + 0.62, cyt, cw - 0.85, 0.5, [(it, 11.5, False, INK)])
    _footer_slide(s, idx)
    return s


def s_summary(prs, idx):
    s = _blank(prs)
    summary_slide(
        s,
        "Summary",
        "Four commitments of the AI Factory",
        [
            (CLAY, "Value-led", "Every model has an owner & measured ROI"),
            (SKY, "Platform-powered", "Reusable assets cut cost & time per use case"),
            (OLIVE, "Governance-assured", "Responsible, compliant, secure by design"),
            (FIG, "People-first", "Talent & culture sustain the advantage"),
        ],
        "The AI Factory turns AI opportunity into trusted, scaled business value.",
    )
    return s


def generate():
    prs = Presentation()
    prs.slide_width = PInches(SW)
    prs.slide_height = PInches(SH)

    builders = [
        s_title,
        s_purpose,
        s_hub_spoke,
        s_loop,
        s_responsibilities,
        s_org,
        s_portfolio,
        s_stack,
        lambda p, i: _chevron_flow(p, i, "MLOps lifecycle", "Predictive AI",
                                   ["Data &\nfeatures", "Train &\ntune", "Validate\n& test", "Registry", "Deploy", "Monitor &\nretrain"],
                                   "Reproducible pipelines · drift triggers retraining · rollback ready"),
        lambda p, i: _chevron_flow(p, i, "LLMOps lifecycle", "Generative AI",
                                   ["Index /\nembeddings", "RAG /\nground", "Prompt &\norchestrate", "Eval:\nquality+safety", "Guardrails", "Serve &\nobserve"],
                                   "Grounded retrieval · automated + human eval · guardrails · cost/quality monitoring"),
        s_gates,
        s_responsible,
        s_kpi,
        s_requirements,
        s_90days,
        s_summary,
    ]
    for i, b in enumerate(builders, 1):
        b(prs, i)

    out = OUTPUT / "Head-of-AI-Factory-Slides.pptx"
    prs.save(out)
    print(f"Created {out} ({len(prs.slides._sldIdLst)} slides)")
    return out


if __name__ == "__main__":
    generate()
