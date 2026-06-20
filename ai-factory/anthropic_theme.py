"""Anthropic official brand theme for python-pptx slide generators.

Colors and typography per Anthropic brand-guidelines skill:
  https://github.com/anthropics/skills/blob/main/skills/brand-guidelines/SKILL.md

- Dark #141413 · Light #faf9f5 · Accents: clay #d97757 · sky #6a9bcc · olive #788c5d
- Headings: Poppins (Arial fallback) · Body: Lora (Georgia fallback)
"""

from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches as PInches, Pt as PPt
from pptx.dml.color import RGBColor

# Slide dimensions (16:9)
SW, SH = 13.333, 7.5

# --- Anthropic brand palette ---
SLATE = RGBColor(0x14, 0x14, 0x13)
CREAM = RGBColor(0xFA, 0xF9, 0xF5)
IVORY_SUBTLE = RGBColor(0xE8, 0xE6, 0xDC)
IVORY_MEDIUM = RGBColor(0xF0, 0xEE, 0xE6)
MID_GRAY = RGBColor(0xB0, 0xAE, 0xA5)
SLATE_MUTED = RGBColor(0x5E, 0x5D, 0x59)
SLATE_SOFT = RGBColor(0x3D, 0x3D, 0x3A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)

# Accent colors (cycle for shapes — one accent per section max on real Anthropic UI)
CLAY = RGBColor(0xD9, 0x77, 0x57)
SKY = RGBColor(0x6A, 0x9B, 0xCC)
OLIVE = RGBColor(0x78, 0x8C, 0x5D)
FIG = RGBColor(0xC4, 0x66, 0x86)
CACTUS = RGBColor(0xBC, 0xD1, 0xCA)
KRAFT = RGBColor(0xE3, 0xDA, 0xCC)

# Backward-compatible aliases (legacy generator variable names)
NAVY = SLATE
BLUE = SKY
GREEN = OLIVE
ORANGE = CLAY
PURPLE = FIG
TEAL = CACTUS
RED = FIG
INK = SLATE
GREY = SLATE_MUTED
LIGHT = IVORY_SUBTLE
BORDER = IVORY_SUBTLE
LINK = MID_GRAY
ON_DARK = CREAM
ON_DARK_MUTED = MID_GRAY
CARD_FILL = WHITE
SLIDE_BG = CREAM
DARK_CARD = SLATE_SOFT

ACCENTS = [CLAY, SKY, OLIVE, FIG, CACTUS, SLATE_SOFT]
ACCENT_CYCLE = ACCENTS

FONT_HEADING = "Poppins"
FONT_BODY = "Lora"
FONT_FALLBACK_HEADING = "Arial"
FONT_FALLBACK_BODY = "Georgia"


def _font_for(size: float, bold: bool) -> str:
    """Headings 14pt+ bold use Poppins; else Lora."""
    if size >= 14 and bold:
        return FONT_HEADING
    if size >= 22:
        return FONT_HEADING
    return FONT_BODY


def _style_run(p, size, bold, color):
    p.font.size = PPt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = _font_for(size, bold)


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def slide_bg(slide, fill=SLIDE_BG):
    """Full-bleed cream background (Anthropic ivory rhythm)."""
    rect(slide, MSO_SHAPE.RECTANGLE, 0, 0, SW, SH, fill=fill)


def rect(slide, shape_type, x, y, w, h, fill=None, line=None, line_w=None):
    sp = slide.shapes.add_shape(shape_type, PInches(x), PInches(y), PInches(w), PInches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = PPt(line_w or 1.0)
    sp.shadow.inherit = False
    return sp


def text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """runs: list of (text, size, bold, color)."""
    tb = slide.shapes.add_textbox(PInches(x), PInches(y), PInches(w), PInches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, (txt, size, bold, color) in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = txt
        p.alignment = align
        _style_run(p, size, bold, color)
    return tb


def set_shape_text(sp, lines, align=PP_ALIGN.CENTER):
    """lines: list of (text, size, bold, color)."""
    tf = sp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i, (txt, size, bold, color) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = txt
        p.alignment = align
        _style_run(p, size, bold, color)


def header(slide, title, kicker=None, dark=True):
    """Slate banner + clay accent — Anthropic editorial header."""
    slide_bg(slide)
    bar_fill = SLATE if dark else IVORY_MEDIUM
    title_color = CREAM if dark else SLATE
    kicker_color = CLAY if dark else SLATE_MUTED
    rect(slide, MSO_SHAPE.RECTANGLE, 0, 0, SW, 1.15, fill=bar_fill)
    rect(slide, MSO_SHAPE.RECTANGLE, 0, 1.15, SW, 0.06, fill=CLAY)
    if kicker:
        text(slide, 0.55, 0.16, 12.2, 0.3, [(kicker.upper(), 11, True, kicker_color)])
        text(slide, 0.55, 0.42, 12.2, 0.6, [(title, 23, True, title_color)])
    else:
        text(slide, 0.55, 0.3, 12.2, 0.7, [(title, 24, True, title_color)], anchor=MSO_ANCHOR.MIDDLE)
    return 1.55


def footer(slide, idx, label="Anthropic-style deck"):
    text(slide, 0.55, 7.05, 9.0, 0.3, [(label, 9, False, SLATE_MUTED)])
    text(slide, 11.8, 7.05, 1.0, 0.3, [(str(idx), 9, True, SLATE_MUTED)], align=PP_ALIGN.RIGHT)


def card(slide, x, y, w, h, accent, heading, sub=None, big=None):
    """White card on cream slide, clay/sky/olive left bar."""
    rect(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, fill=CARD_FILL, line=BORDER, line_w=1.0)
    rect(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, 0.12, h, fill=accent)
    cy = y + 0.18
    if big:
        text(slide, x + 0.25, cy, w - 0.4, 0.9, [(big, 34, True, accent)])
        cy += 0.92
    text(slide, x + 0.28, cy, w - 0.45, 0.5, [(heading, 14, True, INK)])
    if sub:
        text(slide, x + 0.28, cy + 0.42, w - 0.45, h - (cy - y) - 0.5, [(sub, 10.5, False, GREY)])


def title_slide(slide, eyebrow, headline, sublines=None, footer_line=None):
    """Full-bleed slate title (Anthropic dark band)."""
    rect(slide, MSO_SHAPE.RECTANGLE, 0, 0, SW, SH, fill=SLATE)
    rect(slide, MSO_SHAPE.RECTANGLE, 0, 4.65, SW, 0.08, fill=CLAY)
    for i, c in enumerate([CLAY, SKY, OLIVE]):
        rect(slide, MSO_SHAPE.OVAL, 11.35 + (i % 2) * 0.65, 0.75 + (i // 2) * 0.65, 0.35, 0.35, fill=c)
    text(slide, 0.9, 2.15, 11, 0.4, [(eyebrow.upper(), 13, True, CLAY)])
    text(slide, 0.9, 2.6, 11.5, 1.3, [(headline, 44, True, CREAM)])
    yy = 4.05
    for ln in sublines or []:
        text(slide, 0.9, yy, 11.5, 0.55, [(ln, 16, False, ON_DARK_MUTED)])
        yy += 0.52
    if footer_line:
        text(slide, 0.9, 5.55, 11.5, 0.55, [(footer_line, 13, True, CREAM)])


def summary_slide(slide, eyebrow, headline, pillars, closing):
    """Dark summary slide with accent-bordered cards."""
    rect(slide, MSO_SHAPE.RECTANGLE, 0, 0, SW, SH, fill=SLATE)
    text(slide, 0.9, 0.8, 11.5, 0.4, [(eyebrow.upper(), 12, True, CLAY)])
    text(slide, 0.9, 1.15, 11.5, 0.85, [(headline, 28, True, CREAM)])
    cw, gap = 2.85, 0.3
    x0 = (SW - (cw * len(pillars) + gap * (len(pillars) - 1))) / 2
    for i, (acc, head, sub) in enumerate(pillars):
        x = x0 + i * (cw + gap)
        rect(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, 2.45, cw, 2.55, fill=SLATE_SOFT, line=acc, line_w=2)
        badge = rect(slide, MSO_SHAPE.OVAL, x + cw / 2 - 0.3, 2.7, 0.6, 0.6, fill=acc)
        set_shape_text(badge, [(str(i + 1), 18, True, CREAM)])
        text(slide, x + 0.2, 3.5, cw - 0.4, 0.5, [(head, 14, True, CREAM)], align=PP_ALIGN.CENTER)
        text(slide, x + 0.22, 4.0, cw - 0.5, 0.95, [(sub, 10.5, False, ON_DARK_MUTED)], align=PP_ALIGN.CENTER)
    rect(slide, MSO_SHAPE.RECTANGLE, 0.9, 5.55, 11.5, 0.06, fill=CLAY)
    text(slide, 0.9, 5.75, 11.5, 0.75, [(closing, 17, True, CREAM)])


# Legacy aliases for generators that use underscore-prefixed names
_blank = blank
_rect = rect
_text = text
_set_shape_text = set_shape_text
_header = header
_footer = footer
_card = card
